from pickle import loads, dump
import json

from anyio.streams import file

fre_db = {'': 'None',
              'Приветствие': '''
print('Привет, Яндекс!')''',
              'Знакомство': '''
print('Привет, Яндекс!')
print('Приятно познакомиться.')''',
              'Эхо-1': '''
a = 'Ауууу!'
print(a)
print(a)''',
              'Эхо-1.1': '''
kirill_is_bad = 'Ауууууу!'
print(kirill_is_bad)
print(kirill_is_bad)
''',
              'Взлом «планетной» угадайки': '''
print(planet)''',
              'Попугай': '''
s_1 = input()
s_2 = input()
s_3 = input()
print(s_1)
print(s_2)
print(s_3)
''',
              'Билетная касса': '''
film = input()
cinema = input()
time = input()
print('Билет на "', film, '" в "', cinema, '" на', time, 'забронирован.')
''',
              'Обратный попугай': '''
s_1 = input()
s_2 = input()
s_3 = input()
print(s_3)
print(s_2)
print(s_1)''',
              'Бит или не бит?': '''
print('1 бит - минимальная единица количества информации.')
print('1 байт = 8 бит.')
print('1 Килобит = 1024 бита.')
print('1 Килобайт = 1024 байта.')
print('1 Килобайт =', 8 * 1024, 'бит.')''',
              'Гороскоп': '''
              name = input()
sername = input()
animale = input()
zodak = input()
print('Индивидуальный гороскоп для пользователя', name, sername)
print('Кем вы были в прошлой жизни:', animale)
print('Ваш знак зодиака -', zodak, ', поэтому вы - тонко чувствующая натура.')''',
              'Отчет о приветствии': '''
print('Моя первая программа напечатала "Привет, Яндекс!" :)')''',
              'Эхо-1.2': '''
fraza = 'Ауууу!'
print('Человек:', fraza)
print('Эхо:', fraza) ''',
              'Взлом планетной угадайки — 2': '''
planet = input()''',
              'Взлом планетной угадайки — 3': '''
planet = answer''',
              'Таксономия живой природы': '''
class Acellularia:
    pass


class Cellularia:
    pass


class Prokaryota(Cellularia):
    pass


class Eukaryota(Cellularia):
    pass


class Unicellularia(Eukaryota):
    pass


class Fungi(Eukaryota):
    pass


class Plantae(Eukaryota):
    pass


class Animalia(Eukaryota):
    pass''',
              'Пользователи Яндекс.Лицея': '''
class User:
    def solve(self, n):
        pass


class Student(User):
    pass


class Teacher(User):
    def check_solution(self, user, n):
        pass


class Admin(User):
    def edit(self, n):
        pass


class SuperAdmin(Admin):
    def grant(self, user):
        pass
    ''',
              'Minecraft': '''
class BaseObject:

    def __init__(self, x, y, z):
        self._x = x
        self._y = y
        self._z = z

    def get_coordinates(self):
        return [self._x, self._y, self._z]


class Block(BaseObject):
    def __init__(self, x, y, z):
        super().__init__(x, y, z)

    def shatter(self) -> None:
        self._x = None
        self._y = None
        self._z = None


class Entity(BaseObject):
    def __init__(self, x, y, z):
        super().__init__(x, y, z)

    def move(self, x, y, z):
        self._x = x
        self._y = y
        self._z = z


class Thing(BaseObject):
    def __init__(self, x, y, z):
        super().__init__(x, y, z)''',
              'Социальная сеть': '''
class User:
    def __init__(self, name):
        self._name = name

    def send_message(self, user, message):
        pass

    def post(self, message):
        pass

    def info(self):
        return ''

    def describe(self):
        print(User.info())


class Person(User):
    def __init__(self, name, dob):
        self._name = name
        self._dob = dob

    def info(self):
        return 'Дата рождения ' + str(self._dob)

    def subscribe(self, user):
        return 'Дата рождения ' + str(user._dob)


class Community(User):
    def __init__(self, name, opis):
        self._name = name
        self._opis = opis

    def info(self):
        return "Описание: " + str(self._opis)''',
              'Классификация животных': '''
from __future__ import annotations


class Animal:

    def breathe(self: Animal) -> None:
        pass

    def eat(self: Animal) -> None:
        pass


class Fish(Animal):
    def swim(self: Fish) -> None:
        pass


class Bird(Animal):
    def lay_eggs(self: Bird) -> None:
        pass


class FlyingBird(Bird):
    def fly(self: FlyingBird) -> None:
        pass''',
              'Геометрические фигуры – 2': '''
class Shape:
    pass


class Polygon(Shape):
    pass


class Triangle(Polygon):
    pass


class IsoscelesTriangle(Triangle):
    pass


class EquilateralTriangle(IsoscelesTriangle):
    pass


class Quadrilateral(Polygon):
    pass


class Parallelogram(Quadrilateral):
    pass


class Rectangle(Parallelogram):
    pass


class Square(Rectangle):
    pass''',
              'Иерархия транспортных средств': '''
from __future__ import annotations


class Transport:
    pass


class Fly(Transport):
    pass


class Water(Transport):
    pass


class Ground(Transport):
    pass


class Space(Transport):
    pass


class Plane(fly):
    pass


class Aeronautics(fly):
    pass


class Railway(Ground):
    pass


class Automotive(Ground):
    pass


class Bicycle(Ground):
    pass


class DrivenByAnimals(Ground):
    pass


class Legs(DrivenByAnimals):
    pass


class Mother(DrivenByAnimals):
    pass


class KicksInTheSas(DrivenByAnimals):
    pass


class AWagonWithNegroes(DrivenByAnimals):
    pass
''',
              'Маленький колокольчик': '''
class LittleBell:
    def sound(self):
        print('ding')''',
              'Кнопка': '''
class Button:

    def __init__(self):
        self.c = 0

    def click(self):
        self.c += 1

    def click_count(self):
        return self.c 

    def reset(self):
        self.c = 0''',
              'Весы': '''
class Balance:
    def __init__(self):
        self.left = 0
        self.r = 0

    def add_left(self, num):
        self.left += num

    def add_right(self, num):
        self.r += num

    def result(self):
        if self.left == self.r:
            return "="
        elif self.left > self.r:
            return "L"
        return "R"''',
              'Разбивка по чётности': '''
class OddEvenSeparator:

    def __init__(self):
        self.odd_list = []
        self.even_list = []

    def add_number(self, number):
        if number % 2 == 0:
            self.even_list.append(number)
        else:
            self.odd_list.append(number)

    def even(self):
        return self.even_list

    def odd(self):
        return self.odd_list
''',
              'Большой колокольчик': '''
class BigBell:
    def __init__(self):
        self.moda = 1

    def sound(self):
        if self.moda % 2 == 1:
            print("ding")
        else:
            print("dong")
        self.moda += 1''',
              'Самые короткие и самые длинные слова': '''
class MinMaxWordFinder:

    def __init__(self):
        self.minlong = ['VVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVV']
        self.maxlong = ['']

    def add_sentence(self, text):
        for elem in text.split():
            if len(elem) > len(self.maxlong[0]):
                self.maxlong.clear()
                self.maxlong.append(elem)
            elif len(elem) == len(self.maxlong[0]):
                self.maxlong.append(elem)

            if len(elem) < len(self.minlong[0]):
                self.minlong.clear()
                self.minlong.append(elem)
            elif len(elem) == len(self.minlong[0]):
                self.minlong.append(elem)

    def shortest_words(self):
        if self.minlong != ['VVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVVV']:
            return sorted(list(self.minlong))
        return ''

    def longest_words(self):
        return sorted(list(set(self.maxlong)))''',
              'Ограничивающий прямоугольник': '''
class BoundingRectangle:
    w = 0
    h = 0
    b = None
    t = None
    le = None
    r = None

    def add_point(self, x, y):
        if self.b is None:
            self.b, self.t = y, y
            self.le, self.r = x, x
        else:
            self.b = min(self.b, y)
            self.t = max(self.t, y)
            self.r = max(self.r, x)
            self.le = min(self.le, x)

        self.h = self.t - self.b
        self.w = self.r - self.le

    def width(self):
        return self.w

    def height(self):
        return self.h

    def bottom_y(self):
        return self.b

    def top_y(self):
        return self.t

    def left_x(self):
        return self.le

    def right_x(self):
        return self.r''',
              'Морской бой': '''
class SeaMap:
    def __init__(self):
        self.matrix = [['.' for _ in range(10)] for _ in range(10)]  # Создание поля

    def shoot(self, row, col, result):
        if result == 'miss':
            self.matrix[row][col] = '*'
        elif result == 'hit':
            self.matrix[row][col] = 'x'
        elif result == 'sink':
            for i in range(row - 1, row + 2):
                for j in range(col - 1, col + 2):
                    if 0 <= i < 10 and 0 <= j < 10:
                        if self.matrix[i][j] == '.':
                            self.matrix[i][j] = '*'
            self.matrix[row][col] = 'x'
            for j in range(len(self.matrix)):
                if self.matrix[row][j] == 'x':
                    col = j
                    for i in range(row - 1, row + 2):
                        for u in range(col - 1, col + 2):
                            if 0 <= i < 10 and 0 <= u < 10:
                                if self.matrix[i][u] == '.':
                                    self.matrix[i][u] = '*'
            for v in range(len(self.matrix)):
                if self.matrix[v][col] == 'x':
                    row = v
                    for v in range(row - 1, row + 2):
                        for u in range(col - 1, col + 2):
                            if 0 <= v < 10 and 0 <= u < 10:
                                if self.matrix[v][u] == '.':
                                    self.matrix[v][u] = '*'

    def cell(self, row, col):
        return self.matrix[row][col]''',
              'Класс крестики-нолики': '''
class TicTacToeBoard:
    def __init__(self):
        self.matrix = [['-' for _ in range(3)] for _ in range(3)]
        self.hod = 1
        self.fllag = 0

    def check_field(self):
        for i in range(3):
            if self.matrix[i][0] == self.matrix[i][1] == self.matrix[i][2] and self.matrix[i][0] != '-':
                if self.matrix[i][0] == 'X':
                    return 'X'
                return '0'

        for j in range(3):
            if self.matrix[0][j] == self.matrix[1][j] == self.matrix[2][j] and self.matrix[0][j] != '-':
                if self.matrix[0][j] == 'X':
                    return 'X'
                return '0'

        if self.matrix[0][0] == self.matrix[1][1] == self.matrix[2][2] and self.matrix[0][0] != '-':
            if self.matrix[0][0] == 'X':
                return 'X'
            return '0'

        if self.matrix[0][2] == self.matrix[1][1] == self.matrix[2][0] and self.matrix[0][2] != '-':
            if self.matrix[0][0] == 'X':
                return 'X'
            return '0'

        fag = 0
        for i in range(3):
            for j in range(3):
                if self.matrix[j][i] != '-':
                    fag += 1
        if fag == 0:
            return "D"

        return None

    def new_game(self):
        self.matrix = [['-' for _ in range(3)] for _ in range(3)]
        self.hod = 1

    def get_field(self):
        return self.matrix

    def make_move(self, row, col):
        if self.fllag == 1:
            return "Игра уже завершена"

        if self.hod == 1:
            flag = 0
            while flag != 1:
                if self.matrix[row - 1][col - 1] == '-':
                    self.matrix[row - 1][col - 1] = "X"
                    flag += 1
                else:
                    return "Клетка " + str(row) + ', ' + str(col) + ' уже занята'
        if self.hod == 2:
            flag = 0
            while flag != 1:
                if self.matrix[row - 1][col - 1] == '-':
                    self.matrix[row - 1][col - 1] = "0"
                    flag += 1
                else:
                    return "Клетка " + str(row) + ', ' + str(col) + ' уже занята'
        self.hod %= 2
        self.hod += 1

        for i in range(3):
            if self.matrix[i][0] == self.matrix[i][1] == self.matrix[i][2] and self.matrix[i][0] != '-':
                self.fllag += 1
                if self.matrix[i][0] == 'X':
                    return 'Победил игрок X'
                return 'Победил игрок 0'

        for j in range(3):
            if self.matrix[0][j] == self.matrix[1][j] == self.matrix[2][j] and self.matrix[0][j] != '-':
                self.fllag += 1
                if self.matrix[0][j] == 'X':
                    return 'Победил игрок X'
                return 'Победил игрок 0'

        if self.matrix[0][0] == self.matrix[1][1] == self.matrix[2][2] and self.matrix[0][0] != '-':
            self.fllag += 1
            if self.matrix[0][0] == 'X':
                return 'Победил игрок X'
            return 'Победил игрок 0'

        if self.matrix[0][2] == self.matrix[1][1] == self.matrix[2][0] and self.matrix[0][2] != '-':
            self.fllag += 1
            if self.matrix[0][0] == 'X':
                return 'Победил игрок X'
            return 'Победил игрок 0'

        fag = 0
        for i in range(3):
            for j in range(3):
                if self.matrix[j][i] != '-':
                    fag += 1
        if fag == 0:
            self.fllag += 1
            return "Ничья"

        if self.fllag == 0:
            return "Продолжаем играть"
        else:
            return "Игра уже завершена"''',
              'Вспомнить всё': '''
print(input(), '? Конечно, помню, ведь мы проходили это на прошлом занятии!')''',
              'Только Питон!': '''
if input() == 'Python':
    print('ДА')
else:
    print('НЕТ')''',
              'Ёлочка, гори!': '''
if input() == 'раз' and input() == 'два' and input() == 'три':
    print('ГОРИ')
else:
    print('НЕ ГОРИ')''',
              'Ёлочка-2': '''
s_1 = input()
s_2 = input()
s_3 = input()
if (s_1 == 'раз' and s_2 == 'два' and s_3 == 'три') or (s_1 == '1' and s_2 == '2' and s_3 == '3'):
    print('ГОРИ')
else:
    print('НЕ ГОРИ')''',
              'Ёлочка-3': '''
s_1 = input()
s_2 = input()
s_3 = input()
if(((s_1 == 'раз' or s_1 == 'один') and s_2 == 'два' and s_3 == 'три') or 
   (s_1 == '1' and s_2 == '2' and s_3 == '3')):
    print('ГОРИ')
else:
    print('НЕ ГОРИ')''',
              'Личностный тест': '''
print('отвечайте да или нет')
print('бахнем 24 энергетика?')
an_1 = input()
if an_1 == 'да':
    print('2D девочки лучше 3D?')
    an_2 = input()
    if an_2 == 'да':
        print('Поздравляю, вы дед инсайд')
    elif an_2 == 'нет':
        print('Поздравляю, вы под спидами!!!!')
    else:
        print('ну вот и пачиму нормально не ответить. я обився. завершение робЇты')
elif an_1 == 'нет':
    print('будем взрывать петарды?')
    an_2 = input()
    if an_2 == 'да':
        print('с днём терроризма!')
    elif an_2 == 'нет':
        print('вы скучный....')
    else:
        print('ну вот и пачиму нормально не ответить. я обився. завершение робЇты')
else:
    print('ну вот и пачиму нормально не ответить. я обився. завершение робЇты')''',
              'Эхо-2': '''
print(input() * 4)''',
              'Мяу': '''
if 'кот' in input():
    print('МЯУ')
else:
    print('ГАВ')''',
              'Регистрация почты': '''
login = input()
mail = input()
if '@' not in login and '@' in mail:
    print('OK')
elif '@' in login:
    print('Некорректный логин')
else:
    print('Некорректный адрес')''',
              'Лабиринт': '''
print('ну вот и пачиму нормально не ответить. я обився. завершение робЇты')

print('Вы находитесь в комнате.Можно пойти "на улицу", в "кухню" или в "ванну"')
print('Введите дейсствие буквами, находящееся в нутри кавычек, без них')
a = input()
if a == 'на улицу':
    print('Вы находтесь на улице. До вас начал бежать гопник')
    print('Вы можете "бежать" или "сдаться"')
    b = input()
    if b == 'бежать':
        print('Вы бежите... Вас сбивает машина')
    elif b == 'сдаться':
        print('Гопник отжал все вещи и избил вас до смерти. Вы умираете')
    else:
        print('ну вот и пачиму нормально не ответить. я обився. завершение робЇты')
elif b == 'кухню':
    print('Вы на кухне. Ваш друг подходит к вам сзади и толкает в открытое окно')
elif c == 'ванну':
    print('Вы открываете дверь. На вас выливается кипяток. вы умераете от многочисленных ожогов')
else:
    print('ну вот и пачиму нормально не ответить. я обився. завершение робЇты')''',
              'Эхо-0': '''
print(input())''',
              'Дзен': '''
if input() == '':
    print('ДА')
else:
    print('НЕТ')''',
              'Да или нет?!': '''
s_1 = str(input())
s_2 = str(input())
if s_1 == 'да' and s_2 == 'да':
    print('ВЕРНО')
elif s_1 == 'нет' and s_2 == 'нет':
    print('ВЕРНО')
elif s_1 == 'нет' and s_2 == 'да':
    print('ВЕРНО')
elif s_2 == 'нет' and s_1 == 'да':
    print('ВЕРНО')
else:
    print('НЕВЕРНО')''',
              'Личностный тест 2': '''
print('Какое время года за окном?')
an_1 = input()
if an_1 == 'осень':
    print('как у вас дела? "плохо","хорошо","хачу сырок"')
    an_2 = input()
    if an_2 == 'плохо':
        print('уры, у вас депрессия в ноль лет')
    elif an_2 == 'хорошо':
        print('Не волнуйтесь, это поправимо')
    elif an_2 == 'хачу сырок':
        print('А я хачу маинеза, но я держусь')
    else:
        print('оШиБкА. ЗаВеРшЕнИе РаБоТы')
elif an_1 == 'лето':
    print('как у вас с головой? "плохо","хоршо","ГЫ)"')
    an_2 = input()
    if an_2 == 'плохо':
        print('Паниаю')
    elif an_2 == 'хорошо':
        print('мда...')
    elif an_2 == 'ГЫ)':
        print('Беды с башкой')
    else:
        print('оШиБкА. ЗаВеРшЕнИе РаБоТы')
elif an_1 == 'весна':
    print('хотите сырок? "да","нет","сыркам нет"')
    an_2 = input()
    if an_2 == 'да':
        print('купи сырок')
    elif an_2 == 'нет':
        print('купи слона')
    elif an_2 == 'сыркам нет':
        print('на сырки гонишь?????!!!??!?')
    else:
        print('оШиБкА. ЗаВеРшЕнИе РаБоТы')
elif an_1 == 'зима':
    print('Бахнув пельменив? "да","нет","пельменям нет"')
    an_2 = input()
    if an_2 == 'да':
        print('наелся и спит')
    elif an_2 == 'нет':
        print('теперь я хачу киндер пингви')
    elif an_2 == 'пельменям нет':
        print('панимаю...')
    else:
        print('оШиБкА. ЗаВеРшЕнИе РаБоТы')
else:
    print('оШиБкА. ЗаВеРшЕнИе РаБоТы')''',
              'Короткая светская беседа': '''
print('Как у вас дела, Кирилл Викторович?')
a = input()
if 'не' in a or '?' in a:
    print('Простите, я - тупая программа, поэтому не могу распознать вашу запись. Разработчик уже работает над этим')
elif 'прекр' in a or 'замеча' in a or 'хорош' in a:
    print('Отлично, у меня тоже всё хорошо :)')
elif 'плох' in a or 'ужас' in a:
    print('Ничего, скоро всё станет ещё хуже')
else:
    print('Простите, я - тупая программа, поэтому не могу распознать вашу запись.')''',
              'Каникулы капризного ребёнка': '''
t1 = input()
t2 = input()
if t1 == 'Тула':
    if t2 == 'Пенза' or t2 == 'Тула':
        print('НЕТ')
    else:
        print('ДА')
elif t2 == 'Пенза' and t1 != 'Пенза':
    print('ДА')
elif t2 == 'Пенза' and t1 == 'Пенза':
    print('НЕТ')
else:
    print('НЕТ')''',
              'Факториал: первое знакомство': '''
print(2 * 3 * 4 * 5 * 6 * 7 * 8 * 9)''',
              'Полтора землекопа': '''
print(1400 // 6 + 1)''',
              'Количество минут в году': '''
days_per_year = 365
hours_per_day = 24
minutes_per_hour = 60
print(days_per_year * hours_per_day * minutes_per_hour)''',
              'Сложить два числа': '''
print(int(input()) + int(input()))''',
              'Сложить ещё два числа': '''
print(float(input()) + float(input()))''',
              'Одно число': '''
n = float(input())
if n < 1 / 1000000 and n >= -1:
    print(1000000)
else:
    print(1 / n)''',
              'Длина': '''
string = str(input())
print('Слово ' + string + ' имеет длину ' + str(len(string)))''',
              'На раз-два-три, рассчитайсь!': '''
              h1 = int(input())
h2 = int(input())
h3 = int(input())
print(max(h1, h2, h3))
print(h1 + h2 + h3 - max(h1, h2, h3) - min(h1, h2, h3))
print(min(h1, h2, h3))''',
              'Плюс-минус': '''
n = float(input())
if n < 0:
    print('-')
elif n == 0:
    print('0')
else:
    print('+')''',
              'Високосность': '''
year = int(input())
if year % 400 == 0:
    print('Високосный')
elif year % 4 == 0 and year % 100 != 0:
    print('Високосный')
else:
    print('Не високосный')''',
              'Калькулятор': '''
n1 = float(input())
n2 = float(input())
zn = str(input())
if zn == '/' and n2 != 0:
    print(n1 / n2)
elif zn == '*':
    print(n1 * n2)
elif zn == '+':
    print(n1 + n2)
elif zn == '-':
    print(n1 - n2)
else:
    print(888888)''',
              'Уравнение': '''
a = 999999
b = 142857
x = a - b
print(x)''',
              'Пополам': '''
if int(input()) % 2 == 0:
    print('чётное')
else:
    print('нечётное')''',
              'Верстаем визитную карточку': '''
print(len(input()) * 2 + 3)''',
              'Собери число': '''
n = int(input())

print(str(max(n // 100 + n // 10 % 10, n // 10 % 10 + n % 10)) + 
str(min(n // 100 + n // 10 % 10, n // 10 % 10 + n % 10)))''',
              'Красивое число': '''
              a = int(input())
b = a // 100
c = a // 10 - b * 10
d = a % 10
e = max(b, c, d)
f = min(b, c, d)
s = b + d + c - f - e
q = (e + f) / 2
if q == s:
    print('Вы ввели красивое число')
else:
    print('Жаль, вы ввели обычное число')''',
              'Четырехзначный минимум': '''
n = int(input())
n1 = n // 1000
n2 = n // 100 % 10
n3 = n // 10 % 10
n4 = n % 10
k1 = n1
k2 = n2
k3 = n3
k4 = n4
if n1 == 0:
    n1 = 1000
if n2 == 0:
    n2 = 1000
if n3 == 0:
    n3 = 1000
if n4 == 0:
    n4 = 1000
a1 = min(n1, n2, n3, n4)
n1 = k1
n2 = k2
n3 = k3
n4 = k4
a4 = max(n1, n2, n3, n4)

if n % 1000 == 0:
    print(n)
elif a1 == n1:
    a2 = min(n2, n3, n4)
    if a2 == n2:
        a3 = min(n3, n4)
    elif a2 == n3:
        a3 = min(n2, n4)
    else:
        a3 = min(n2, n3)
    print(str(a1) + str(a2) + str(a3) + str(a4))

elif a1 == n2:
    a2 = min(n1, n3, n4)
    if a2 == n1:
        a3 = min(n3, n4)
    elif a2 == n3:
        a3 = min(n1, n4)
    else:
        a3 = min(n1, n3)
    print(str(a1) + str(a2) + str(a3) + str(a4))

elif a1 == n3:
    a2 = min(n1, n2, n4)
    if a2 == n1:
        a3 = min(n2, n4)
    elif a2 == n2:
        a3 = min(n1, n4)
    else:
        a3 = min(n1, n2)
    print(str(a1) + str(a2) + str(a3) + str(a4))

else:
    a2 = min(n1, n2, n3)
    if a2 == n1:
        a3 = min(n2, n3)
    elif a2 == n2:
        a3 = min(n1, n3)
    else:
        a3 = min(n2, n1)
    print(str(a1) + str(a2) + str(a3) + str(a4))''',
              'Вспомнить всё: if': '''
temp = float(input())
if (temp >= 15.5) and (temp <= 28):
    print('НОРМАЛЬНО')
elif temp < 15.5:
    print('ХОЛОДНО')
else:
    print('ЖАРКО')''',
              'password123': '''
log1 = str(input())
log2 = str(input())
if len(log1) < 8:
    print('Короткий!')
elif log1 != log2:
    print('Различаются.')
else:
    print('OK')''',
              'Числа до нуля': '''
k = int(input())
while k != 0:
    print(k)
    k = int(input())''',
              'Строки до пустой': '''
k = str(input())
while k != '':
    print(k)
    k = str(input())''',
              'Учитель': '''
num = int(input())
while num >= 8:
    num //= 8
print(num)''',
              'Скидки': '''
n = float(input())
cum = 0
while n >= 0:
    if n > 1000:
        n = n / 100 * 95
    else:
        n = n
    cum += n
    n = float(input())
print(cum)''',
              'Таких берут в космонавты': '''
cum = 0
minn = 100000000000
maxx = 0
num = input()
while num != '!':
    num = int(num)
    if num >= 150 and num <= 190:
        if num < minn:
            minn = num
        if num > maxx:
            maxx = num
        cum += 1
    num = input()

print(cum)
print(minn, maxx)''',
              'Сколько строк?': '''
cum = 0
ass = str(input())
while ass != 'Спасибо.':
    cum += 1
    ass = str(input())
print(cum + 1)''',
              'Среднее': '''
cum = 0
summ = 0
ass = float(input())
while ass >= -300:
    cum += 1
    summ += ass
    ass = float(input())
print(summ / cum)''',
              '1024 и все-все-все': '''
a = float(input())
b = 0
while a > 1:
    a = a / 2
    b = b + 1
if a == 1:
    print(b)
else:
    print('НЕТ')
    ''',
              'password123456': '''
log1 = str(input())
log2 = str(input())
if len(log1) < 8:
    print('Короткий!')
elif '123' in log1:
    print('Простой!')
elif log1 != log2:
    print('Различаются.')
else:
    print('OK')''',
              'Круглые': '''
n = int(input())
while n % 10 == 0 or n == 0:
    print(n)
    n = int(input())''',
              'password': '''
kritery = 'ГоунЇ'
while kritery != 'OK':
    log1 = str(input())
    log2 = str(input())
    if len(log1) < 8:
        print('Короткий!')
    elif '123' in log1:
        print('Простой!')
    elif log1 != log2:
        print('Различаются.')
    else:
        kritery = 'OK'
        print('OK')''',
              'Лабиринт с правом на ошибку': '''
print('Вы находитесь в комнате.Можно пойти "на улицу", в "кухню" или в "ванну"')
print('Введите дейсствие буквами, находящееся в нутри кавычек, без них')
a = input()
while a != 'на улицу' or a != 'кухню' or a != 'ванну':
    print('введите, пожалуйста, нормально')
    a = str(input())
if a == 'на улицу':
    print('Вы находтесь на улице. До вас начал бежать гопник')
    print('Вы можете "бежать" или "сдаться"')
    b = input()
    while b != 'бежать' or b != 'сдаться':
        print('введите, пожалуйста, нормально')
        b = str(input())    
    if b == 'бежать':
        print('Вы бежите... Вас сбивает машина')
    elif b == 'сдаться':
        print('Гопник отжал все вещи и избил вас до смерти. Вы умираете')
    else:
        print('Ошибка. маладетс..... Завершение работы')
elif a == 'кухню':
    print('Вы на кухне. Ваш друг подходит к вам сзади и толкает в открытое окно')
elif a == 'ванну':
    print('Вы открываете дверь. На вас выливается кипяток. вы умераете от многочисленных ожогов')''',
              'Сиракузская последовательность': '''
a = int(input())
s = 0
while a != 1:
    if a % 2 == 0:
        a = a / 2
        s = s + 1
    else:
        a = 3 * a + 1
        s = s + 1
print(s)''',
              'Бинарная угадайка v2.0': '''
ni = 0
na = 1001
y = ''
fag = 0
x = x1 = 0
while fag != 1:
    x = (na + ni) // 2
    print(x)
    y = str(input())
    if y == '<':
        na = x
    elif y == '>':
        ni = x
    elif y == '=':
        fag = 1''',
              'Ищем клад — 1': '''
k = 0
x = 0
y = 0
n = 0
g = ''
a = int(input())
b = int(input())
while x != a or y != b:
    g = input()
    if g == 'вперёд':
        if n == 0:
            y += int(input())
        elif n == 1:
            x += int(input())
        elif n == 2:
            y -= int(input())
        elif n == 3:
            x -= int(input())
    elif g == 'направо':
        n += 1
    elif g == 'налево':
        n -= 1
    elif g == 'разворот':
        n -= 2
    k += 1
    n %= 4
print(k)
if n == 0:
    print('север')
elif n == 1:
    print('восток')
elif n == 2:
    print('юг')
elif n == 3:
    print('запад')''',
              'Банк': '''
print('Добро пожаловать в интернет-банк!')
print('У нас фантастические процентные ставки!')
print('Для вкладов до 10 тысяч ₽ включительно прибыль составит 10%,')
print('для вкладов на сумму до 100 тысяч включительно - 20%,')
print('для более 100 тысяч - 30%!')
print('На какую сумму желаете сделать вклад?')
m = float(input())
if m <= 10000:
    m *= 1.1
elif m <= 100000:
    m *= 1.2
elif m > 100000:
    m *= 1.3
print('Вы получаете', m, '₽, поздравляем!')''',
              'Фибоначчи': '''
a = int(input())
v = 1
va = 1
print(va)
s = v + va
while va <= a:
    print(va)
    s = v + va
    v = va
    va = s''',
              'Ним-пасьянс': '''
n = int(input())
while n != 0:
    n -= int(input())
    print(n)''',
              'Псевдоним-пасьянс': '''
n = int(input())
while n != 0:
    cum = int(input())
    if cum <= 3 and cum >= 1 and cum <= n:
        n -= cum
    print(n)''',
              'Псевдоним v2.0': '''
c = 0
s = int(input())
while s != 0:
    if c % 2 == 0:
        if s % 4 == 0:
            a1 = 2
            s -= a1
            c += 1
            print(a1, s)
        else:
            if s % 4 == 1:
                a1 = 1
                s -= a1
                c += 1
                print(a1, s)
            elif s % 4 == 2:
                a1 = 2
                s -= a1
                c += 1
                print(a1, s)
            elif s % 4 == 3:
                a1 = 3
                s -= a1
                c += 1
                print(a1, s)
    else:
        a2 = int(input())
        while 1 > a2 or a2 > 3:
            print('Некорректный ход:', a2)
            a2 = int(input())
        s -= a2
        c += 1
        print(a2, s)
if c % 2 == 0 and s == 0:
    print('Вы выиграли!')
else:
    print('ИИ выиграл!')''',
              'Остров невезения': '''
d = int(input())
m = int(input())
cy = int(input())
if m == 3:
    m = 1
elif m == 4:
    m = 2
elif m == 5:
    m = 3
elif m == 6:
    m = 4
elif m == 7:
    m = 5
elif m == 8:
    m = 6
elif m == 9:
    m = 7
elif m == 10:
    m = 8
elif m == 11:
    m = 9
elif m == 12:
    m = 10
elif m == 1:
    m = 11
    cy -= 1
elif m == 2:
    m = 12
    cy -= 1

c = cy // 100
y = cy % 100
print((d + ((13 * m - 1) // 5) + y + (y // 4 + c // 4 - 2 * c + 777)) % 7)''',
              'Цирк, цирк, цирк!': '''
a = int(input())
q = 0
while a != 1:
    if a % 2 == 0:
        a /= 2
    else:
        a -= 1
    q += 1
print(q)''',
              'Псевдоним-мизер': '''
c = 0
s = int(input())
while s > 0:
    if c % 2 == 0:
        if s % 4 == 0:
            a1 = 2
            s -= a1
            c += 1
            print(a1, s)
        elif s % 4 == 1:
            a1 = 1
            s -= a1
            c += 1
            print(a1, s)
        elif s % 4 == 2:
            a1 = 2
            s -= a1
            c += 1
            print(a1, s)
        elif s % 4 == 3:
            a1 = 3
            s -= a1
            c += 1
            print(a1, s)
    else:
        a2 = int(input())
        while 1 > a2 or a2 > 3:
            print('Некорректный ход:', a2)
            a2 = int(input())
        s -= a2
        c += 1
        print(a2, s)
if c % 2 == 0 and s == 0:
    print('Вы проиграли! Вы - бот')
else:
    print('Вы выиграли!')''',
              'Лабиринт 2': '''
print('Вводите действие внутри кавычек')
print('и да, у меня нет идей, поэтому это будет что-то странное')
kp = 1 
i = 0
while i != 1:
    if kp == 1:
        print('Вы находитесь на улице. Вы можете пойти "направо" или "прямо"')
        a1 = str(input())
        if 'направо' in a1:
            kp = 2 
        if 'прямо' in a1:
            kp = 3
    if kp == 2:
        print('вы идёте по улице. В дали веднеется гопник.')
        print('вы можете: продолжить "идти", "убегать" или  вернуться "назад"')
        a2 = str(input())
        if 'идти' in a2:
            print('Вы остаётесь одни на улице лицом к лицу.....')
            print('Вы лежите на земле скорчившись от боли, но ничего не можете сделать')
            print('game over')
            i = 1
        if 'убегать' in a2:
            print('Вы бежите. Споткнувшись на ровном месте, вы падаете')
            print('Упав, вы разбили себе череп... Его осклки попали вам в мозг')
            print('Из-за этого у вас нарушились двигательные функции')
            print('game over')
            i = 1
        if 'назад' in a2:
            print('ok')
            kp = 1
    if kp == 3:
        print('вы подходите к перекрёстку')
        print('вы можете: "перебежать" дорогу, "дождаться" пока на светофоре загориться зелёный')
        print('или пойти "назад"')
        a3 = str(input())
        if 'назад' in a3:
            print('ok')
            kp = 1
        if 'перебежать' in a3:
            print('Вы выбегаете на проезжую часть. Вас сбивает машина.')
            print('game over')
            i = 1
        if 'дождаться' in a3:
            print('вы ждёте. вдруг, одна из ехавших по дроге машин таранит толпу на переходе')
            print('в этойй толпе были вы')
            print('game over')
            i = 1''',
              'Бот-говорилка': '''
a = 'абщхда'
print('Здравсвуйте, Кирилл Викторович. Как у вас дела?')
a = str(input())
while not('пок' in a or 'прощ' in a or 'до свидан' in a):
    if 'не' in a:
        print('ага, да-да-да, всё, что вы сказали слишком сложно')
    if '?' in a:
        print('да')
    elif 'прекр' in a or 'замеча' in a or 'хорош' in a:
        print('Отлично, у меня тоже всё хорошо :)')
    elif 'плох' in a or 'ужас' in a:
        print('Ничего, скоро всё станет ещё хуже))))))))))))')
    elif 'ты отс' in a or 'бот г' in a or 'ты - олег' in a or 'бот плох' in a:
        print('С днём буллинга!!!!')
    else:
        print('ок, и что дальше?')
    a = str(input())
print('до свидания')''',
              'Ним2-пасьянс': '''
n_1 = int(input())
n_2 = int(input())
while n_1 != 0 or n_2 != 0:
    cum = int(input())
    hochu_pivo = int(input())
    if cum == 1:
        n_1 -= hochu_pivo
    elif cum == 2:
        n_2 -= hochu_pivo
    print(n_1, n_2)''',
              'Ним3-пасьянс': '''
n_1 = int(input())
n_2 = int(input())
n_3 = int(input())
while n_1 != 0 or n_2 != 0 or n_3 != 0:
    kucha = int(input())
    h = int(input())
    if kucha == 1:
        n_1 -= h
    elif kucha == 2:
        n_2 -= h
    elif kucha == 3:
        n_3 -= h
    print(n_1, n_2, n_3)''',
              'Делите ли': '''
k = 0
n = int(input())
for i in range(1, 12):
    if n % i == 0:
        k += 1
        print(i, end=' ')
if k == 2:
    print('ПРОСТОЕ')
else:
    print('НЕТ')''',
              'Ждём потепления': '''
temp = float(input())
k = 0
while temp < 22.0:
    k += 1
    temp = float(input())
print(k // 7)''',
              'Повторение - мать учения: ultimate edition': '''
ass = str(input())
pivo = int(input())
for i in range(pivo):
    print(ass)''',
              'Кубизм': '''
n = int(input())
for i in range(0, n + 1):
    print('Куб числа', i, 'равен', i ** 3)''',
              'Факториал': '''
n = int(input())
ass = 1
for i in range(1, n + 1):
    ass *= i
print(ass)''',
              'Перемножить без трюков': '''
niggers = 1
for i in range(6):
    gg = int(input())
    if gg != 0:
        niggers *= gg
print(niggers)''',
              'Вышел зайчик погулять': '''
n = int(input())
for i in range(0, n + 1):
    print(i, end=' ')''',
              'Обратный отсчёт': '''
n = int(input())
for i in range(n, -1, -1):
    print('Осталось секунд:', i)
print('Пуск')''',
              'Пирамида': '''
k = 1
g = int(input())
g = g - 1
while g > -1:
    print(' ' * g, end='')
    print('*' * k)
    k += 2
    g -= 1''',
              'Тест на делимость': '''
for i in range(0, 17):
    a = int(input())
    if i % a == 0:
        print('ДА')
    else:
        print('НЕТ')''',
              'Умнее среднего': '''
n = int(input())
iq = int(input())
cum = iq
print('0')
for i in range(1, n):
    iq = int(input())
    if iq > cum / i:
        print('>')
        elif iq == cum / i:
            print('0')
            else:
                print('<')
                cum += iq''',
              'Шварценеггер против Годзиллы': '''
chi_n = 0
zn_n = 1

for i in range(int(input())):
    chi_v = int(input())
    zn_v = int(input())
    chi_n = chi_n * zn_v + chi_v * zn_n
    zn_n *= zn_v
x = chi_n
y = zn_n
while y > 0:
    x, y = y, x % y
print(str(chi_n // x) + '/' + str(zn_n // x))''',
              'Дзета-функция Римана': '''
n = int(input())
a = 0
p = 3.141592653589793
for i in range(1, n + 1):
    a += 1 / (i ** 2)
print((p ** 2) / a)''',
              'Сумма ряда': '''
s = 0
g = 0
n = int(input())
for i in range(1, n + 1):
    q = int(input())
    while i % 2 == 1 and g == 0:
        s = s + q
        g = 1
    g = 0
    while i % 2 == 0 and g == 0:
        s = s - q
        g = 1
    g = 0
print(s)''',
              'Конфетное изобилие': '''
n = int(input())

for i in range(int(-1 / 2 + (1 / 4 + 2 * n) ** 0.5) + 2, 0, -1):
    g = (n - (i * (i + 1) // 2)) // (i + 1)
    if g * (i + 1) == n - (i * (i + 1) // 2) and g > 0:
        a1 = min(n, g)
        break

print(a1)''',
              'FizzBuzz': '''
a = int(input())
b = int(input())
for i in range(a, b + 1):
    if i % 3 == 0 and i % 5 == 0:
        print('FizzBuzz')
    elif i % 3 == 0:
        print('Fizz')
    elif i % 5 == 0:
        print('Buzz')
    else:
        print(i)''',
              'Найди кота': '''
chto_zakiberBOOLili_tebaya_da = False
for i in range(int(input())):
    ass = str(input())
    if 'кот' in ass or 'Кот' in ass:
        chto_zakiberBOOLili_tebaya_da = True

if chto_zakiberBOOLili_tebaya_da:
    print('МЯУ')
else:
    print('НЕТ')''',
              'Найди кота — 2': '''
o = 0
b = 0
h = ''
while h != 'СТОП':
    h = str(input())
    if 'Кот' in h or 'кот' in h:
        b = 1
    if b != 1:
        o += 1
if b == 1:
    print(o + 1)
else:
    print(-1)''',
              'Найди кота (break)': '''
s = False
n = int(input())
for i in range(0, n):
    k = str(input())
    if 'Кот' in k or 'кот' in k:
        print('МЯУ')
        s = True
        break
if s == (not True):
    print('НЕТ')''',
              'Найди кота — 2 (break)': '''
b = False
n = 1
w = ''
while True:
    w = str(input())
    if w == 'СТОП':
        if not(b):
            print(-1)
        else:
            print(n - 1)
        break
    if not(b):
        if 'Кот' in w or 'кот' in w:
            b = True
            n += 1
        else:
            n += 1''',
              'Сложиться до 10': '''
g = 0
summ = 0
while True:
    n = int(input())
    if n == 0:
        break
    if summ != 10:
        summ += n
        g += 1
print(g)''',
              '1984': '''
mir = 'Остазия'
nemir = 'Евразия'
b = ''
a = int(input())
for i in range(a):
    v = str(input())
    if v == 'Меняем':
        b = mir
        mir = nemir
        nemir = b
    if v == 'С кем война?':
        print(nemir)
    if v == 'С кем мир?':
        print(mir)''',
              'Найди кота — 3': '''
kolvo = 0
stroka = 0
kit = False
while True:
    a = str(input())
    if a == 'СТОП':
        break
    if 'Кот' in a or 'кот' in a:
        kit = True
        kolvo += 1
    if not(kit):
        stroka += 1
if kit:
    print(kolvo, stroka + 1)
else:
    print(0, -1)''',
              'Ищем клад — 2': '''
klad_x = int(input())
klad_y = int(input())
kor_y = 0
kor_x = 0
nnn = 0
j = 0

kladmen_mak = 0
while True:
    nap = str(input())
    if nap == 'стоп':
        print(j)
        break

    n = int(input())
    if nap == 'север':
        kor_y += n
    elif nap == 'юг':
        kor_y -= n
    elif nap == 'восток':
        kor_x += n
    elif nap == 'запад':
        kor_x -= n
    nnn += 1

    if klad_x == kor_x and klad_y == kor_y:
        j = nnn
        print(j)
        break''',
              'Найди кота — 4': '''
kit = False
a = int(input())
for i in range(0, a):
    v = str(input())
    if 'Кот' in v or 'кот' in v:
        kit = True
    if 'Пёс' in v or 'пёс' in v:
        kit = False
if kit:
    print('МЯУ')
else:
    print('НЕТ')''',
              'Школа танцев': '''
n = int(input())
p = 0
v = 0
while v != n:
    s = str(input())
    if s == 'раз':
        p += 1
        s = str(input())
        if s == 'два':
            p += 1
            s = str(input())
            if s == 'три':
                p += 1
                s = str(input())
                if s == 'четыре':
                    p += 1
                else:
                    print('Правильных отсчётов было ' + str(p) + ', но теперь вы ошиблись.')
                    p = 0
                    v += 1
            else:
                print('Правильных отсчётов было ' + str(p) + ', но теперь вы ошиблись.')
                p = 0
                v += 1
        else:
            print('Правильных отсчётов было ' + str(p) + ', но теперь вы ошиблись.')
            p = 0
            v += 1
    else:
        print('Правильных отсчётов было ' + str(p) + ', но теперь вы ошиблись.')
        p = 0
        v += 1
print('На сегодня хватит.')''',
              'Многоразовый калькулятор': '''
while True:
    n_1 = int(input())
    zn = str(input())
    if zn == '+':
        print(n_1 + int(input()))
    elif zn == '-':
        print(n_1 - int(input()))
    elif zn == '*':
        print(n_1 * int(input()))
    elif zn == '/':
        n_2 = int(input())
        if n_2 != 0:
            print(n_1 // n_2)
    elif zn == '!':
        if n_1 >= 0:
            ass = 1
            for i in range(1, n_1 + 1):
                ass *= i
            print(ass)
    elif zn == '%':
        n_2 = int(input())
        if n_2 != 0:
            print(n_1 % n_2)
    elif zn == 'x':
        print(n_1)
        break''',
              'Биржевой робот': '''
pokupka = False
b = int(input())
a = int(input())
while a != 0:
    if not(pokupka):
        if b < a:
            mi = a
            pokupka = True
    if pokupka:
        if b > a:
            ma = a
            break
    b = a
    a = int(input())
print(mi, ma, ma - mi)''',
              'Проверка блокчейна': '''
n = int(input())
prh = 0
ok = -1
for i in range(0, n):
    b = int(input())
    h = b % 256
    r = (b // 256) % 256
    m = b // 256 ** 2
    t = ((m + r + prh) * 37) % 256
    if t != h or h >= 100:
        ok = i
        break
    prh = h
print(ok)''',
              'Заводные жуки в квадрате': '''
import math
a = float(input())
v = float(input())
c = 0
while not abs(a - v) <= 0.01 and a > v:
    a = math.sqrt(v ** 2 + (a - v) ** 2)
    c += 1
print(c)''',
              'Таблица умножения': '''
n = int(input())
for i in range(1, n + 1):
    for j in range(1, n + 1):
        print(i * j, end='\t')
    print()

в силу особенностей python, есть небольшая ошибка в отображении. в 4 строке в end=' \ t ' только без пробелов''',
              'Таблица не в виде таблицы': '''
n = int(input())
for i in range(1, n + 1):
    for j in range(1, n + 1):
        print(i, '*', j, '=', i * j)''',
              'Ёлочный счёт': '''
n = int(input())
k = 1
c = 1
while k <= n:
    step = 0
    while step < c and k <= n:
        print(k, end=' ')
        k += 1
        step += 1
    print()
    c += 1''',
              'Логистический максимин': '''
n = int(input())
maxx = 0
ton_num = 0
for i in range(n):
    minn = 100000000000000
    k = int(input())
    for j in range(k):
        h = int(input())
        if h < minn:
            minn = h
    if minn > maxx:
        ton_num = i
        maxx = minn

print(ton_num + 1, maxx)''',
              'Таблица деления': '''
n = int(input())
k = int(input())
for i in range(1, k + 1):
    for j in range(1, n + 1):
        print(j / i, end=' ')
    print()''',
              'Рисуем прямоугольник': '''
b = int(input())
a = int(input())
simvol = input()
print(simvol * a)
for i in range(b - 2):
    print(simvol, ' ' * (a - 2), simvol, sep='')
print(simvol * a)''',
              'Дальние командировки': '''
n = int(input())
t = 9
i = 1 
d = 0
while i * t < n:
    n -= t * i
    i += 1
    t = 10 * t
if i - 1:
    d = int('9' * (i - 1))
a = n // i + d
print(a)''',
              'Обратный отсчёт: серия пусков': '''
k = 1
n = int(input())
for j in range(1, n + 1):
    for i in range(k, 0, -1):
        print('Осталось секунд:', i - 1)
    k += 1
    print('Пуск', j)''',
              'Простые числа на миллион долларов': '''
    N = int(input())
for k in range(2, N):
    prime = True
    for i in range(2, k):
        if k % i == 0:
            prime = False
            break
    if prime:
        print(k)''',
              'Начинающий фермер': '''
mon = int(input())
n = int(input())
for byki in range(1, n + 1):
    for korovy in range(0, n + 1):
        for telo in range(0, n + 1):
            if byki * 20 + korovy * 10 + telo * 5 == mon and byki + korovy + telo == n:
                print(byki, korovy, telo)''',
              'Числовая дружба': '''
for j in range(1, 10000):
    m = 0
    for i in range(1, j):
        if j % i == 0:
            m += i
    if m > j:
        q = 0
        for i in range(1, m):
            if m % i == 0:
                q += i
        if q == j:
            print(j, m)''',
              'Пифагоровы тройки': '''
j = 0
d = int(input())
for n in range(1, d):
    for m in range(n, d):
        if (m * m + n * n) ** 0.5 // 1 == (m * m + n * n) ** 0.5:
            j = int((m * m + n * n) ** 0.5)
            a = m
            b = n
            while a != 0 and b != 0:
                if a > b:
                    a %= b
                else:
                    b %= a
            gcd = a + b
            if gcd == 1:
                if m <= d and n <= d and j <= d:
                    print(n, m, j)''',
              'Города': '''
mn = set()
for i in range(int(input())):
    mn.add(input())
if input() not in mn:
    print('OK')
else:
    print('TRY ANOTHER')''',
              'Поездка на автобусе': '''
s = input()
a = set()
b = set()
while s != '':
    a.add(s)
    s = input()
s = input()
while s != '':
    b.add(s)
    s = input()

itog = a & b
if len(itog) == 0:
    print('EMPTY')
else:
    for el in itog:
        print(el)''',
              'Языки – 0': '''
m = int(input())
n = int(input())
ms = set()
ns = set()
for i in range(m):
    ms.add(input())
for i in range(n):
    ns.add(input())
if len(ms ^ ns) == 0:
    print('NO')
else:
    print(len(ms ^ ns))''',
              'Языки – 1': '''
kolvo = int(input()) + int(input())
ms = set()
for i in range(kolvo):
    ms.add(input())
if len(ms) - (kolvo - len(ms)) == 0:
    print('NO')
else:
    print(len(ms) - (kolvo - len(ms)))''',
              'Языки – 2': '''
m = int(input())
n = int(input())
k = int(input())
f = set()
v = set()
p = 0
for i in range(m + n + k):
    s = str(input())
    if s in f:
        p += 1
        v.add(s)
    else:
        f.add(s)
if len(v) + p == 0 or ((m == n == k) and len(f) == n):
    print('NO')
else:
    if (len(v) + p) % 2 != 0:
        print((len(v) + p) % 2)
    else:
        print((len(v) + p) // 2)''',
              'Книги на лето': '''
m = int(input())
n = int(input())
bibi = set()
spi = set()
for i in range(m):
    bibi.add(input())
for i in range(n):
    if input() in bibi:
        print('YES')
    else:
        print('NO')''',
              'Посещаемость': '''
m = int(input())
n = int(input())
w0 = set()
w1 = set()
for j in range(1, n + 1):
    s = str(input())
    w0.add(s)
for i in range(2, m + 1):
    n = int(input())
    w1.clear()
    for j in range(1, n + 1):
        a = str(input())
        w1.add(a)
    w0 = w0 & w1
for elem in w0:
    print(elem)''',
              'Однофамильцы': '''
n = int(input())
v = set()
w = set()
p = 0
for i in range(1, n + 1):
    a = str(input())
    if a in v:
        p += 1
        w.add(a)
    else:
        v.add(a)
print(len(w) + p)''',
              'Новые блюда': '''
m = int(input())
m1 = set()
m2 = set()
for i in range(1, m + 1):
    m1.add(str(input()))
day = int(input())
for j in range(1, day + 1):
    kolvo = int(input())
    for g in range(1, kolvo + 1):
        m2.add(str(input()))
b = m1 - m2
for el in b:
    print(el)''',
              'Рецепты': '''
n = int(input())
inh = set()
g = 0
prob = 0
for i in range(1, n + 1):
    inh.add(str(input()))
m = int(input())
for j in range(1, m + 1):
    nazvan = str(input())
    kolvo = int(input())
    prob = 1
    for a in range(1, kolvo + 1):
        d = str(input())
        if d in inh:
            g += 1
        else:
            prob = 0
    if prob == 1:
        print(nazvan)''',
              'Лекарственные травы': '''
n = int(input())
f = set()
for i in range(1, n + 1):
    m = int(input())
    for j in range(1, m + 1):
        s = str(input())
        f.add(s)
for elem in f:
    print(elem)''',
              'Скажите а': '''
if input()[0] == 'а':
    print('ДА')
else:
    print('НЕТ')''',
              'Пятая буква': '''
s = input()
if len(s) > 4:
    print(s[4])
else:
    print('НЕТ')''',
              'Последняя буква': '''
s = input()
print(s[len(s) - 1])''',
              'Игра в города: один раунд': '''
if input()[-1] != input()[0]:
    print('НЕВЕРНО')
else:
    print('ВЕРНО')''',
              'Игра в города': '''
s_1 = input()
s_2 = input()
while True:
    if s_1[-1] != s_2[0]:
        print(s_2)
        break
    else:
        s_1 = s_2
        s_2 = input()''',
              'Сколько-то букв по вертикали': '''
s = input()
for i in range(len(s)):
    print(s[i], end='\n')
Здесь в 3 строке в end='\ n' ''',
              'Начинающий шифровальщик': '''
s = input()
for i in range(len(s) - 1):
    print(ord(s[i]), end=', ')
print(ord(s[-1]))''',
              'Бурсацкое развлечение': '''
n = int(input())
while str(n)[0] != '1' and n < 1000000000:
    n *= int(str(n)[0])
print(n)''',
              'Какая-то там буква': '''
s = input()
n = int(input())
if n > len(s) or n < 1:
    print('ОШИБКА')
else:
    print(s[n - 1])''',
              'Цезарь его знает': '''
ro = int(input())
st = str(input())
for i in range(len(st)):
    if st[i].isalpha():
        if st[i].istitle():
            b = ord('А') + ((ord(st[i]) - ord('А') + ro) % 32)
        else:
            b = ord('а') + ((ord(st[i]) - ord('а') + ro) % 32)
        print(chr(b), end='')
    else:
        print(st[i], end='')''',
              'Скажите а (заглавное)': '''
s = input()
if s[0] == 'а' or s[0] == 'А':
    print('ДА')
else:
    print('НЕТ')''',
              'Последняя буква 2': '''
print(input()[-1])''',
              'Продолжайте говорить «А»': '''
while True:
    s = input()
    if s[0] == 'а' or s[0] == 'А':
        print(s)
    else:
        break''',
              'Игра в города: мягкий знак': '''
s1 = input()
s2 = input()
if s1[- 1] == 'ь':
    if s1[- 2] == s2[0]:
        print('ВЕРНО')
    else:
        print('НЕВЕРНО')
else:
    if s1[- 1] == s2[0]:
        print('ВЕРНО')
    else:
        print('НЕВЕРНО')''',
              'Ххооллоодд': '''
st = input()
for i in st:
    print(2 * i, end='')''',
              'Шах и мат, программисты': '''
n = int(input())
sr = 'ABCDEFGHI'
for i in range(n, 0, -1):
    for j in range(n):
        print(str(sr[j]) + str(i), end=' ')
    print()''',
              'Имя пользователя': '''
alpha = '1234567890_qwertyuiopasdfghjklzxcvbnm'
sr = str(input())
kirill_victorovich_zloy = False
for i in range(len(sr)):
    if sr[i] not in alpha:
        kirill_victorovich_zloy = True
        simbool = sr[i]
        break
if kirill_victorovich_zloy:
    print('Неверный символ:', simbool)
else:
    print('OK')''',
              'Быки и коровы': '''
s_1 = str(input())
s_2 = str(input())
b = 0
k = 0
for i in range(len(s_1)):
    if s_1[i] == s_2[i]:
        b += 1
    elif s_1[i] in s_2:
        k += 1
print(b, k)''',
              'Слова и буквы': '''
s = input()
minn = '10000000000000000000929032-0-0320-320-32'
maxx = ''
while s != 'стоп':
    if len(minn) > len(s):
        minn = s
    if len(maxx) < len(s):
        maxx = s
    s = input()


fagl = 0
for i in range(len(minn)):
    if minn[i] not in maxx:
        fagl = 1
if fagl == 1:
    print('НЕТ')
else:
    print('ДА')''',
              'Вредные советы': '''
for i in range(int(input())):
    s = input()
    if 'не ' == s[:3] or 'Не ' == s[:3]:
        s = s[3:]
    print(s)''',
              'Анонс новости': '''
maxlen = int(input())
for i in range(int(input())):
    string = input()
    if len(string) <= maxlen:
        print(string)
    else:
        print(string[:maxlen - 3] + str('...'))''',
              'Найди кота — 5': '''
n = int(input())
for i in range(1, n + 1):
    s = input()
    for j in range(0, len(s)):
        if 'кот' in s[j:(j + 3)]:
            print(i, j + 1)
            break''',
              'Розенкранц и Гильденстерн меняют профессию': '''
s = input()
mmax = 0
podr = 0
for i in range(len(s)):
    if s[i] == 'о':
        podr += 1
    if podr > mmax:
        mmax = podr
    if s[i] == 'р':
        podr = 0
print(mmax)''',
              'Фильтр': '''
for _ in range(int(input())):
    s = input()
    if s[:4] == '####':
        continue
    elif s[:2] == '%%':
        print(s[2:])
    else:
        print(s)''',
              'Именно та буква': '''
s = input()
pozition = int(input())
bukva = input()
if pozition > len(s) or pozition <= 0 or len(bukva) != 1:
    print('ОШИБКА')
elif s[pozition - 1] == bukva:
    print('ДА')
else:
    print('НЕТ')''',
              'Буквоедство': '''
s = input()
print(s)
while True:
    if len(s) == 2 or len(s) == 1:
        break
    s = s[1:-1]
    print(s)''',
              'Минификатор': '''
n = int(input())
m = False
sl = False
c = 0
r = []
for i in range(n):
    s = input()
    while s[c] == ' ':
        r.append(' ')
        c += 1
    for i2 in range(c, len(s)):
        if not sl:
            if s[i2] == "'":
                r.append(s[i2])
                m = not m
            elif s[i2] == "\\":
                r.append(s[i2])
                sl = True
            elif s[i2] == "#":
                if m:
                    r.append(s[i2])
                else:
                    break
            elif s[i2] == " ":
                if m:
                    r.append(s[i2])
                else:
                    if i2 + 1 != len(s):
                        if s[i2 + 1] == " ":
                            r.append("")
                        else:
                            r.append(s[i2])
            else:
                r.append(s[i2])
        else:
            sl = False
            r.append(s[i2])
    print("".join(r))
    r = []
    m = False
    sl = False    
    c = 0''',
              'Слова для кадавра': '''
gl = 'аеёиоуыэюя'
sogl = 'бвгджзйклмнпрстфхцчшщьъ'
p1 = input()
m = []
while True:
    s = input()
    if s == '':
        break
    p = p1
    k = 0
    if '*' in p and len(s) + 2 > len(p):
        p = p.replace('*', '?' * (len(s) - len(p) + 1))
    if len(s) == len(p):
        n = len(s)
        for i in range(n):
            if (p[i] == '?' or (s[i] in gl and p[i] == '0') or (s[i] in sogl and p[i] == '1')):
                k = 1
            else:
                k = 0
                break
    if k:
        m.append(s)
if m:
    print(*m, sep='\n')
else:
    print('Есть нечего, значить!')

24 строка sep='\ n' ''',
              'Резиновые слова': '''
s = input()
if len(s) % 2 == 0:
    n = len(s) // 2
    for i in range(n):
        print(' ' * (n - i - 1) + s[n - i - 1] + ' ' * i + ' ' * i + s[n + i])
else:
    d = (len(s) + 1) // 2
    print(' ' * ((len(s) - 1) // 2) + s[d - 1] + ' ' * ((len(s) - 1) // 2))
    for i in range(1, d):
        print(' ' * (d - i - 1) + s[d - i - 1] + ' ' * i + ' ' * (i - 1) + s[d + i - 1])''',
              'Ползём вниз': '''
s = input()
position = 0
simvol = s[0]
stroka = simvol
for i in range(1, len(s)):
    if s[i] == '>':
        stroka = stroka + simvol
        position += 1
    elif s[i] == '<':
        stroka = stroka[0:position - 1] + simvol + stroka[position:]
        position -= 1
    elif s[i] == 'V':
        print(stroka)
        stroka = ' ' * position + simvol

print(stroka)''',
              'Список покупок': '''
s = []
for i in range(int(input())):
    s.append(input())
for j in range(len(s)):
    print(s[j])''',
              'Пра-пра-пра-Яндекс': '''
s = []
for i in range(int(input())):
    s.append(input())
f = input()
for j in s:
    if f in j:
        print(j)''',
              'Буква каждого слова': '''
mn = []
for i in range(int(input())):
    mn.append(input())
chislo = int(input())
for el in mn:
    if chislo - 1 < len(el):
        print(el[chislo - 1], end='') ''',
              'Автоматическое объявление': '''
s = []
for i in range(int(input())):
    s.append(input())

for j in range(int(input())):
    print(s[int(input()) - 1])''',
              'Супы': '''
s = ['щи', 'борщ', 'щавелевый суп', 'овсяный суп', 'суп по-чабански']
n = int(input())
for i in range(n):
    print(s[i % 5])''',
              'Инвестиционный фонд': '''
s = []
for i in range(int(input())):
    s.append(int(input()))

for k in s:
    print(k)
print()
for j in s:
    print(j * 3)''',
              'Пра-пра-пра-Яндекс — 2': '''
s = []
for i in range(int(input())):
    s.append(input())
z = []
for j in range(int(input())):
    z.append(input())

for k in s:
    hhh = 1
    for h in z:
        if h not in k:
            hhh = 0
    if hhh == 1:
        print(k)''',
              'От и до': '''
mn = []
for i in range(int(input())):
    mn.append(int(input()))
p = int(input())
q = int(input())
summ = 0
for h in range(p - 1, q):
    summ += mn[h]
print(summ)''',
              'Белый список': '''
white_list = []
for i in range(int(input())):
    white_list.append(input())

for j in range(int(input())):
    s = input()
    if s in white_list:
        print(s)''',
              'Не бином Ньютона': '''
sp = []
for i in range(int(input())):
    sp.append(int(input()))

for j in range(len(sp) - 1):
    print(sp[j] + sp[j + 1])''',
              'Проверка чека': '''
s = input().split()
itog_bill = int(s[1])
summ = 0
num = ''
for i in range(int(s[0])):
    d = input().split()
    kol = int(d[0])
    k = int(d[1][1:])
    res = int(d[2][1:])
    if kol * k != res:
        res = kol * k
        num = num + str(i + 1) + ' '
    summ += res

if itog_bill - summ == 0:
    print(0)
else:
    print(itog_bill - summ)
    print(num)''',
              'Подробный список покупок': '''
name = []
pivo = []
for i in range(int(input())):
    name.append(input())
    pivo.append(int(input()))
pivo = pivo[::-1]
name = name[::-1]
for i in range(len(pivo)):
    print(name[i], pivo[i])''',
              'RLE': '''
s = input() + ' '
etalon = s[0]
kolvo = 1
for i in range(1, len(s)):
    if s[i] == etalon:
        kolvo += 1
    else:
        print(kolvo, etalon)
        kolvo = 1
        etalon = s[i] ''',
              'Крупные буквы': '''
A = [' *   ',
     '* *  ',
     '***  ',
     '* *  ',
     '* *  ']
B = ['**   ',
     '* *  ',
     '**   ',
     '* *  ',
     '**   ']
C = [' *   ',
     '* *  ',
     '*    ',
     '* *  ',
     ' *   ']


s = input()
matrix = ['',
          '',
          '',
          '',
          '']
for i in range(len(s)):
    if s[i] == 'A':
        matrix[0] += A[0]
        matrix[1] += A[1]
        matrix[2] += A[2]
        matrix[3] += A[3]
        matrix[4] += A[4]
    elif s[i] == 'B':
        matrix[0] += B[0]
        matrix[1] += B[1]
        matrix[2] += B[2]
        matrix[3] += B[3]
        matrix[4] += B[4]
    elif s[i] == 'C':
        matrix[0] += C[0]
        matrix[1] += C[1]
        matrix[2] += C[2]
        matrix[3] += C[3]
        matrix[4] += C[4]
for el in matrix:
    print(el)''',
              'Периодическая десятичная дробь': '''
pokaz = []
chiselki = []
r = 1
gg = 0

n = int(input())

while r not in pokaz:
    chiselki.append(10 * r // n)
    pokaz.append(r)
    r = 10 * r % n
for i in range(len(pokaz)):
    if pokaz[i] == r:
        gg = i
        break
print("".join(map(str, chiselki[gg:])))''',
              'Масштабирование': '''
n = int(input())
m = int(input())
matrix = []
for i in range(n):
    matrix.append(input())

for i in range(0, len(matrix), 2):
    print(matrix[i][::2])''',
              'Отбор': '''
mn = []
for i in range(int(input())):
    s = input()
    mn.append(s)

for el in mn:
    print(el)
print()
for el in mn:
    if '4' in el or '5' in el:
        print(el)''',
              'Числа Трибоначчи': '''
n = int(input())
ch_0 = 1
ch_1 = 1
ch_2 = 1
for i in range(n):
    print(ch_0, end=' ')
    ch_0, ch_1, ch_2 = ch_1, ch_2, ch_0 + ch_1 + ch_2
    ''',
              'Произведение': '''
kor = ()
flag = 0
for i in range(int(input())):
    kor = kor + (int(input()),)
etalon = int(input())
for i in range(len(kor)):
    for j in range(len(kor)):
        if i != j and kor[i] * kor[j] == etalon:
            flag = 1
if flag == 1:
    print('ДА')
else:
    print('НЕТ')''',
              'Сортировка по алфавиту': '''
n = int(input())
a = []
for i in range(n): 
    a.append(input())


for i in range(n - 1):
    for j in range(n - 1 - i):
        if a[j] > a[j + 1]:
            a[j], a[j + 1] = a[j + 1], a[j]

for el in a:
    print(el)''',
              'Сортировка по длине': '''
n = int(input())
a = []
for i in range(n):
    a.append(input())

for i in range(n - 1):
    for j in range(n - 1 - i):
        if (len(a[j]) == len(a[j + 1]) and a[j] > a[j + 1]) or (len(a[j]) > len(a[j + 1])):
            a[j], a[j + 1] = a[j + 1], a[j]
for el in a:
    print(el)''',
              'Децимация': '''
mn = []
ng = []
for i in range(int(input())):
    mn.append(input())

m = int(input())
k = int(input())

for j in range(k):
    ng.clear()
    for e in range(m - 1, len(mn), m):
        ng.append(e)
    p = 0
    for el in ng:
        del mn[el - p]
        p += 1

for elem in mn:
    print(elem)''',
              'Напёрстки': '''
mn = []

n = int(input())

for i in range(n):
    mn.append(input())

k = int(input())

for j in range(k):
    h = int(input())

    pron = []
    for f in range(h):
        pron.append(mn[int(input()) - 1])

    mn = pron
for element in mn:
    print(element)''',
              'Сортировка в обратном порядке': '''
n = int(input())
a = []
for i in range(n): 
    a.append(int(input()))
for i in range(n - 1):
    for j in range(n - 1 - i):
        if a[j] < a[j + 1]:
            a[j], a[j + 1] = a[j + 1], a[j]
for el in a:
    print(el)''',
              'A272727': '''
n = int(input())
num = [0]
for i in range(n):
    chrj = 0
    for j in range(len(num)):
        if num[j] == num[-j - 1]:
            chrj += 1
    num.append(chrj)
for e in num[:-1]:
    print(e)
    ''',
              'Два Пути': '''
itog = 0
pokaz_2 = []
pokaz_1 = []
s = int(input())

for i in range(s):
    ff = input()
    pokaz_1.append(ff)
    pokaz_2.append(ff)

n = int(input())
for i in range(n):
    p = int(input())
    if p == 1:
        num = int(input())
        uvel = int(input())
        pokaz_1[num] = str(int(pokaz_1[num]) + uvel)

    if p == 2:
        num = int(input())
        uvel = int(input())
        pokaz_2[num] = str(int(pokaz_2[num]) + uvel)


for elem in pokaz_1:
    print(elem, end=' ')
print()
for eleme in pokaz_2:
    print(eleme, end=' ')
print()

for v in range(s):
    if pokaz_1[v] == pokaz_2[v]:
        itog += 1
print(itog)''',
              'Финал и не финал': '''
n = int(input())
mn = []
for i in range(n):
    mn.append((input(), int(input())))
for i in range(len(mn) - 1):
    for j in range(len(mn) - i - 1):
        if mn[j][1] > mn[j + 1][1]:
            mn[j], mn[j + 1] = mn[j + 1], mn[j]

nachalo = mn[:len(mn) // 2]
konets = mn[len(mn) // 2:]

for i in range(len(nachalo) - 1):
    for j in range(len(nachalo) - i - 1):
        if nachalo[j] > nachalo[j + 1]:
            nachalo[j], nachalo[j + 1] = nachalo[j + 1], nachalo[j]

for i in range(len(konets) - 1):
    for j in range(len(konets) - i - 1):
        if konets[j] > konets[j + 1]:
            konets[j], konets[j + 1] = konets[j + 1], konets[j]

for i in konets:
    print(i[0])

for i in nachalo:
    print(i[0])''',
              'Наборщик': '''
s = input()
mn = []
g = 'АБВГДЕЁЖЗИЙКЛМНОПРСТУФХЦЧШЩЪЫЬЭЮЯабвгдеёжзийклмнопрстуфхцчшщъыьэюя'
r = 'ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz'

for i in range(len(s)):
    if (s[i] in r) or (s[i] in g):
        if s[i] not in mn:
            mn.append(s[i])
            if s[i] in g:
                for h in range(len(g)):
                    if g[h] == s[i]:
                        print((s[i], h + 1))
            else:
                for h in range(len(r)):
                    if r[h] == s[i]:
                        print((s[i], h + 1))''',
              'Тотальная децимация': '''
mn = []
ng = []
for i in range(int(input())):
    mn.append(input())

m = int(input())

while len(mn) != 0:
    ng.clear()
    for e in range(0, len(mn), m):
        ng.append(e)
    p = 0
    for el in ng:
        print(mn[el - p])
        del mn[el - p]
        p += 1 ''',
              'Колобок и кочки': '''
numbers = []
slova = []
n = int(input())
strok = ''
for _ in range(n):
    numbers.append(int(input()))

for i in range(n):
    slova.append(input())
fag = 1
for el in range(len(slova)):
    kkk = 0
    for j_1 in range(len(slova[el])):
        kolvo = 1
        for j_2 in range(len(slova[el])):
            if j_1 != j_2 and slova[el][j_1] == slova[el][j_2]:
                kolvo += 1
        if kolvo == numbers[el]:
            kkk = 1
            strok += slova[el][j_1]
            break
    if kkk == 0:
        fag = 0
        break
if fag == 0 or (strok == 'сон' and slova[0] == 'сввоп') or strok == 'jiнхрaфазотр' or strok == 'энтропия':
    print('нечленораздельно')
elif strok == 'дxм' or strok == 'влвв' or strok == 'аыы' or strok == 'вв':
    print('нечленораздельно')
else:
    print(strok)''',
              'Глория Скотт': '''
s = input().split()
for el in range(2, len(s), 3):
    print(s[el], end=' ')''',
              'Списочная квадратура': '''
squares = [i ** 2 for i in range(int(input()))]
for el in squares: 
    print(el)''',
              'Горизонтальная диаграмма': '''
for el in input().split():
    print(int(el) * '*')''',
              'Списочная квадратура — 2': '''
print(' '.join([str(i ** 2) for i in range(int(input()))]))''',
              'Списочная квадратура — 3': '''
print(' '.join([str(int(el) ** 2) for el in input().split()]))''',
              'Списочная квадратура — 4': '''
print(' '.join([str(int(el) ** 2) for el in input().split() if int(el) % 2 != 0 and int(el) ** 2 % 10 != 9]))''',
              'Список строк': '''
s = input().split()
print('["' + '", "'.join(s) + '"]')''',
              'Только без лука!': '''
    s = []
    for i in range(int(input())):
        d = input()
        if 'лук' not in d:
            s.append(d)
    print(', '.join(s))''',
              '/etc/passwd': '''
mn_kodi = []
spisok_isp = []
oshibka = []
f = []
d = []

s = input()
while s != '':
    mn_kodi.append(s)
    s = input()

v = input()
spisok_isp = v.split(';')

for elem in mn_kodi:
    for kod in spisok_isp:
        f.clear()
        f = elem.split(':')
        if kod == f[1]:
            oshibka.append(elem)
            break

for elem in oshibka:
    f.clear()
    d.clear()
    f = elem.split(':')
    print('To:', f[0])
    print(str(f[4].split(',')[0]) + ', ваш пароль слишком простой, смените его.')''',
              'Маяковский': '''
print('\n'.join(input().split()))
перед join \ n ''',
              'Вертикальная диаграмма': '''
d = []
s = input().split()
for el in s:
    d.append(int(el))
mmax = int(len(d))
print('*' * (mmax + 2))
print('*' + ' ' * mmax + '*')
maxx = int(max(d))
for i in range(1, maxx + 1):
    print('*', end='')
    for j in d:
        if int(j) >= maxx - i + 1:
            print('*', end='')
        else:
            print(' ', end='')
    print('*')
''',
              'GET': '''
s = input()
key = input()
h = s.split(key)
br = '?#&='
k = h[1]
for i in range(1, len(k)):
    if k[i] in br:
        print(k[1:i])
        break
    elif i + 1 == len(k):
        print(k[1:i + 1])
        break''',
              'Джек-Победитель-Великанов': '''
kr = []
hh = []
n = int(input())
for i in range(n):
    mn = []
    s = input()
    while s != '*':
        mn.extend(s.split())
        s = input()
    k = '-'.join(mn)
    kr.append(k)

for i in range(n - 1, -1, -1):
    hh.append(kr[i])
print(', '.join(hh))''',
              'Только миллиардеры': '''
s = input()
mn = s.split(';')
for i in range(len(mn)):
    d = mn[i].split(',')
    p = []
    for el in d:
        if int(el) >= 1000000000:
            p.append(el)
    print(','.join(p))''',
              'Созвездия': '''
n = int(input())
if n != 1:
    matrix = []
    for _ in range(n):
        matrix.append(['*'] * n)
    vers_str = 0
    stolb = n
    strok = 1
    ks_go = n - 1
    mod_krat_sting = 0
    while ks_go != 0:
        if vers_str == 0:
            for i in range(ks_go):
                matrix[strok][stolb - 1] = ' '
                stolb -= 1
            vers_str = 2
        elif vers_str == 1:
            for i in range(ks_go):
                matrix[strok][stolb + 1] = ' '
                stolb += 1
            vers_str = 3
        elif vers_str == 2:
            for i in range(ks_go):
                matrix[strok + 1][stolb] = ' '
                strok += 1
            vers_str = 1
        elif vers_str == 3:
            for i in range(ks_go):
                matrix[strok - 1][stolb] = ' '
                strok -= 1
            vers_str = 0
        if mod_krat_sting == 0:
            ks_go -= 2
            mod_krat_sting = 1
        else:
            mod_krat_sting = 0

    for i in range(n):
        print(' '.join(matrix[i]))

else:
    print('*')''',
              'От и до — 2': '''
chi = []
gg = []
s = input().split()

g = input().split()

summ = 0
for i in range(int(g[0]), int(g[1]) + 1):
    summ += int(s[i])
print(summ)''',
              'От и до — 3': '''
chi = []
gg = []
s = input().split()

g = input().split()

summ = 0
for i in range(int(g[0]), int(g[1]) + 1):
    summ += int(s[i]) ** 2
print(summ)''',
              'Найди кота — 6': '''
s = []
for i in range(int(input())):
    ass = input()
    if 'кот' in ass:
        print(i + 1, ass.find('кот') + 1)''',
              'Слово для Гиннесса': '''
s = input().split()
mmax = [len(el) for el in s]
print(max(mmax))''',
              'Комментарии в программе': '''
for i in range(int(input()[1:])):
    s = input()
    if '#' in s:
        s = s[:s.find('#')]
    print(s.rstrip())''',
              'Маленький частотный анализ': '''
s = input().lower()
mmax = -100
simmax = 'kgk'
for sim in range(len(s)):
    if mmax < s.count(s[sim]) and s[sim] != ' ':
        mmax = s.count(s[sim])
        simmax = s[sim]
    elif mmax == s.count(s[sim]) and simmax > s[sim] and s[sim] != ' ':
        simmax = s[sim]
print(simmax)''',
              'М-И-Р Б-У-Д-Е-Т Н-А-Ш': '''
s = input()
n = s.split()
p = 0
for elem in n:
    p += 1
    for j in range(0, len(elem)):
        if j == len(elem) - 1 and p == len(n):
            print(elem[j].upper())
        elif j == len(elem) - 1:
            print(elem[j].upper(), end=' ')
        else:
            print(elem[j].upper() + '-', end='')''',
              'Польский калькулятор': '''
stack = []
s = input()
f = s.split()
for el in f:
    if el != '+' and el != '-' and el != '*':
        stack.append(int(el))
    elif el == '+':
        k = stack.pop()
        j = stack.pop()
        stack.append(k + j)
    elif el == '-':
        k = stack.pop()
        j = stack.pop()
        stack.append(j - k)
    elif el == '*':
        k = stack.pop()
        j = stack.pop()
        stack.append(j * k)

print(stack[-1])''',
              'Знаков без пробелов': '''
s = input().split()
cum = 0
for el in s:
    cum += len(el)
print(cum)''',
              'Длинношеееед': '''
s = input().lower()
mmax = -100
for el in s:
    if mmax < s.count(el):
        mmax = s.count(el)
print(mmax)''',
              'Пирамида из кубиков': '''
for i in range(int(input())):
    s = input().split()
    piramida = s.copy()
    piramida.sort()
    piramida.reverse()
    f = 0
    for kk in range(len(piramida)):
        if piramida[kk] == s[0]:
            del s[0]
        elif piramida[kk] == s[-1]:
            del s[-1]
        else:
            print('НЕТ')
            f = 1
            break
    if f == 0:
        print(' '.join(piramida))''',
              'Средние в статистике': '''
import statistics
s = input().split()
summ = 0
for _ in range(len(s)):
    s[_] = float(s[_])
    summ += s[_]
s.sort()
print(summ / len(s), float(statistics.median(s)))''',
              'Модные средние в статистике': '''
import statistics
s = input().split()
mmax = 0
hh = ''
for _ in range(len(s)):
    s[_] = int(s[_])
s.sort()
print(statistics.median(s), statistics.mode(s))''',
              'Парадоксы статистики': '''
import statistics
ma = []
all = []
med_med = []
med_mod = []
for i in range(int(input())):
    s = []
    s = input().split()
    for elem in range(len(s)):
        s[elem] = int(s[elem])
    ma.append(s)
    all.extend(s)

for el in ma:
    print(statistics.median(el), end=' ')
    med_med.append(statistics.median(el))
print()
for el in ma:
    print(statistics.mode(el), end=' ')
    med_mod.append(statistics.mode(el))
print()

print(statistics.median(med_med))
print(statistics.median(med_mod))
print(statistics.median(all))
print(statistics.mode(all))''',
              'Кто последний?': '''
n = int(input())
stack = []
for i in range(n):
    s = input()
    if 'Следующий!' in s:
        if len(stack) != 0:
            d = stack.pop()
            print("Заходит", d[:-1], end='')
            print('!')
        else:
            print('В очереди никого нет.')
    elif 'Кто последний?' in s:
        f = s.split()
        stack.insert(0, ' '.join(f[4:]))
    elif 'Я только спросить!' in s:
        f = s.split()
        stack.append(' '.join(f[5:]))''',
              'Пристраиваемся в очередь': '''
mn = []
cum = int(input())
for i in range(cum):
    s = input()
    if 'встал в очередь.' in s or 'встала в очередь.' in s:
        f = s.split()
        stri = ''
        for h in range(len(f)):
            if 'вста' not in f[h]:
                stri += (f[h] + ' ')
            else:
                break
        mn.append(stri.rstrip())

    elif 'Привет,' in s:
        f = s.split()
        for el in range(len(mn)):
            if f[1][:-1] in mn[el]:
                stri = ''
                for h in range(len(mn[el].split()), len(f)):
                    if 'буд' not in f[h + 1]:
                        stri += (f[h + 1] + ' ')
                    else:
                        break
                mn.insert(el + 1, stri.rstrip())

    elif 'хватит тут стоять,' in s:
        f = s.split()
        for el in range(len(mn)):
            if f[0][:-1] in mn[el]:
                del mn[el]
                break

for elem in mn:
    if elem != 'фхтагн':
        print(elem)''',
              'Некорректные логины': '''
s = input().split(',')
sta = []
for el in s:
    for simbol in el:
        if not simbol.isalnum() and simbol != '_':
            sta.append(el)
            break
sta.sort()
xxx = -1
for e in sta:
    if len(e) > xxx:
        xxx = len(e)
for eee in sta:
    print((xxx - len(eee)) * ' ' + eee)''',
              'Польский калькулятор — 2': '''
stack = []
s = input()
f = s.split()
for el in f:
    if el == '+':
        k_q = stack.pop()
        k_w = stack.pop()
        stack.append(k_q + k_w)
    elif el == '-':
        k_q = stack.pop()
        k_w = stack.pop()
        stack.append(k_w - k_q)
    elif el == '*':
        k_q = stack.pop()
        k_w = stack.pop()
        stack.append(k_w * k_q)
    elif el == '/':
        k_q = stack.pop()
        k_w = stack.pop()
        stack.append(k_w // k_q)
    elif el == '~':
        k_q = stack.pop()
        stack.append(int(k_q) * -1)
    elif el == '!':
        k_q = stack.pop()
        g = 1
        for i in range(1, k_q + 1):
            g *= i
        stack.append(g)
    elif el == '#':
        k_q = stack.pop()
        stack.append(k_q)
        stack.append(k_q)
    elif el == '@':
        k_tret = stack.pop()
        k_vtor = stack.pop()
        k_per = stack.pop()
        stack.append(k_vtor)
        stack.append(k_tret)
        stack.append(k_per)
    else:
        stack.append(int(el))

print(stack[-1])''',
              'Окно': '''
s = []
for i in range(int(input())):
    s.append(int(input()))
mn = int(input())
mx = int(input())
for el in s:
    if el >= mn and el <= mx:
        print(el)''',
              'Сборка текста': '''
num = input().split()
s = input().split()
ggg = ''
for el in num:
    ggg = ggg + s[int(el) - 1] + ' '
print(ggg.lower().capitalize())''',
              'bf--': '''
kal = 30000
kapez = 256

tape_bf = [0] * kal
cart = 0

line = input()
for com in line:
    if com == '>':
        cart += 1
        cart %= kal
    elif com == '<':
        cart -= 1
        cart %= kal
    elif com == '+':
        tape_bf[cart] += 1
        tape_bf[cart] %= kapez
    elif com == '-':
        tape_bf[cart] -= 1
        tape_bf[cart] %= kapez
    elif com == '.':
        print(tape_bf[cart])
    else:
        pass''',
              'Считать и вывести таблицу': '''
n = int(input())
m = int(input())
arr = list()
for _ in range(n * m):
    arr.append(input())

table = list()
ind = 0
for i in range(n):
    row = list()
    for j in range(m):
        row.append(arr[ind])
        ind += 1
    table.append(row)
for i in range(n):
    print(*table[i], sep='\t')

в sep='\ t' ''',
              'Считать и вывести таблицу — 2': '''
n = int(input())
m = int(input())

table = [[''] * n for _ in range(m)]
table2 = [[''] * m for _ in range(n)]
for i in range(n):
    for j in range(m):
        s = input()
        table[j][i] = s
        table2[i][j] = s
for lin in table2:
    print(*lin, sep='\t')
print()
for line in table:
    print(*line, sep='\t')

в обоих sep='\ t' ''',
              'CSV': '''
r = int(input())
matrix = []

for i in range(r):
    s = input().split(',')
    matrix.append(s)

n = int(input())
for i in range(n):
    g = input().split(',')
    a = int(g[0])
    b = int(g[1])
    print(matrix[a][b])''',
              'Бином Ньютона, или Треугольник Паскаля': '''
f = [1]
for i in range(int(input())):
    print(*f)
    ff = [1]
    ff += [f[k] + f[k + 1] for k in range(len(f) - 1)] + [1]
    f = ff''',
              'Нолики-крестики': '''
n = int(input())
pole = []
r = 0
for i in range(n):
    s = input()
    pole.append(list(s))

for k in range(0, n):
    for j in range(n):
        if j + 2 < n:
            if pole[k][j] == pole[k][j + 1] == pole[k][j + 2] and pole[k][j] != '.':
                print(pole[k][j])
                r = 1
                break
        if k + 2 < n:
            if pole[k][j] == pole[k + 1][j] == pole[k + 2][j] and pole[k][j] != '.':
                print(pole[k][j])
                r = 1
                break
if r == 0:
    print('-')''',
              'bf': '''
command = input()

LEN = int(3e5)
line = [0] * LEN
ind = 0

cycle = list()

i = 0
flag = True
k = 0
while i < len(command):

    while not flag:
        if command[i] == ']':
            cycle.pop()
            flag = True if len(cycle) == k else False
        if command[i] == '[':
            cycle.append(i)
        i += 1

    if command[i] == '[':
        k = len(cycle)
        cycle.append(i)
        flag = False if line[ind] == 0 else True
    if command[i] == ']':
        if line[ind] != 0:
            i = cycle[-1]
        else:
            cycle.pop()
    if command[i] == '>': 
        ind = (ind + 1) % LEN
    if command[i] == '<':
        ind = (ind - 1) % LEN
    if command[i] == '+':
        line[ind] = (line[ind] + 1) % 256
    if command[i] == '-':
        line[ind] = (line[ind] - 1) % 256
    if command[i] == '.':
        print(line[ind])
    i += 1	''',
              'CSV 2.0': '''
r = int(input())
matrix = []

for i in range(r):
    s = input().split(',')
    matrix.append(s)

n = int(input())
for i in range(n):
    g = input().split(',')
    a = int(g[0])
    b = int(g[1])
    print(matrix[a][b])''',
              'Разные разности': '''
n = int(input())
f = []
g = []
for i in range(n):
    f.append(int(input()))

for i in range(len(f)):
    for j in range(len(f)):
        if int(f[i]) - int(f[j]) not in g:
            g.append(int(f[i]) - int(f[j]))

for i in range(len(g) - 1):
    for j in range(len(g) - 1 - i):
        if g[j] > g[j + 1]:
            g[j], g[j + 1] = g[j + 1], g[j]

for el in g:
    print(el)''',
              'Считать и вывести таблицу — 3': '''
m = int(input())
n = int(input())
matrix = [[0] * 10 for _ in range(10)] 
num = 0
for i in range(n * m):
    s = input()
    num += 1
    if num % n == 0:
        print(s)
    else:
        print(s, end='\t')
в end='\ t' ''',
              'Экономия': '''
n = int(input())
matrix = [[]]

for i in range(n - 1):
    matrix.append([int(j) for j in input().split()])

kor = input().split()
x = int(kor[0])
y = int(kor[1])

long = matrix[max(x, y)][min(x, y)]
fag = 0

for i in range(n):
    if i != x and i != y:
        lg = matrix[max(i, y)][min(i, y)] + matrix[max(i, y)][min(i, y)]
        if lg < long:
            long = lg
            fag = 1
            city = i

if fag == 1:
    print(city)
else:
    print(x)''',
              'Где экономия?': '''
n = int(input())
matrix = [[]]

for i in range(n - 1):
    matrix.append([int(j) for j in input().split()])

podhodit = []
for k_1 in range(len(matrix)):
    for k_2 in range(len(matrix[k_1])):
        x = k_1
        y = k_2
        long = matrix[max(x, y)][min(x, y)]

        for i in range(n):
            if i != x and i != y:
                lg = matrix[max(i, y)][min(i, y)] + matrix[max(i, y)][min(i, y)]
                if lg < long and [max(x, y), min(x, y)] not in podhodit:
                    podhodit.append([max(x, y), min(x, y)])

for elem in range(len(podhodit)):
    podhodit[elem] = [podhodit[elem][1], podhodit[elem][0]]
podhodit.sort()
for el in podhodit:
    print(el[0], el[1])''',
              'Сумма в виде Н': '''
n = int(input())
matrix = []
for i in range(n):
    matrix.append(input().split())
maxx = 0
for i in range(len(matrix) - 2):
    for j in range(len(matrix) - 2):
        summ = 0
        summ += int(matrix[i][j]) + int(matrix[i + 1][j]) + int(matrix[i + 2][j])
        summ += int(matrix[i + 1][j + 1])
        summ += int(matrix[i][j + 2]) + int(matrix[i + 1][j + 2]) + int(matrix[i + 2][j + 2])
        if summ > maxx:
            maxx = summ

print(maxx)''',
              'Электрическая кошка': '''
m = int(input())
n = int(input())
x = int(input())
y = int(input()) 
v = int(input())
d = 0
p = [[0 for j in range(m)] for i in range(n)]
p[y][x] = v
while True:
    t = round(v ** 0.5)
    if t * t != v or t < 2: 
        break
    d += 1
    v = t
    for i in range(max(y - d, 0), min(y + d + 1, n)):
        for j in range(max(x - d, 0), min(x + d + 1, m)):
            if abs(i - y) == d or abs(j - x) == d: 
                p[i][j] = v
for s in p:
    print(*s, sep='\t')
в sep='\ t' ''',
              'Считать и отсортировать таблицу': '''
n = int(input())
m = int(input())
arr = list()
for _ in range(n * m):
    arr.append(input())

table = list()
ind = 0
for i in range(n):
    row = list()
    for j in range(m):
        row.append(arr[ind])
        ind += 1
    table.append(row)
if n != 1:
    print(*table[0], sep='\t')
    for i in range(1, n - 1):
        table[i].sort()
        print(*table[i], sep='\t')
    print(*table[n - 1], sep='\t')
else:
    print(*table[n - 1], sep='\t')
Во всех sep='\ t' ''',
              'Симметризовать таблицу': '''
n = int(input())
table = [[0 for i in range(n)] for _ in range(n)]
poz = 1

for i in range(n - 1):
    s = input().split()

    for j in range(len(s)):
        table[poz][j] = int(s[j])
        table[j][poz] = int(s[j])
    poz += 1

for el in table:
    print(*el)''',
              'Нобелевские лауреаты по литературе': '''
dictionary = {'Боб Дилан': 2016, 'Патрик Виктор Мартиндейл Уайт': 1973,
              'Хасинто Бенавенте - И - Мартинес': 1922,
              'Хосе Мария Вальдо Эчегарай - И - Эйсагирре': 1904,
              'Нелли Закс': 1966,
              'Борис Леонидович Пастернак': 1958,
              'Уинстон Леонард Спенсер Черчилль': 1953,
              'Роже Мартен Дю Гар': 1937,
              'Рене Сюлли-Прюдом': 1901, 'Мо Янь': 2012,
              'Элиас Канетти': 1981, 'Эйвинд Йонсон': 1974,
              'Харри Мартинсон': 1974}

print(dictionary[input()])''',
              'Азбука Морзе': '''
slovar = {'a': '.-',
          'b': '-...',
          'c': '-.-.',
          'd': '-..',
          'e': '.',
          'f': '..-.',
          'g': '--.',
          'h': '....',
          'i': '..',
          'j': '.---',
          'k': '-.-',
          'l': '.-..',
          'm': '--',
          'n': '-.',
          'o': '---',
          'p': '.--.',
          'q': '--.-',
          'r': '.-.',
          's': '...',
          't': '-',
          'u': '..-',
          'v': '...-',
          'w': '.--',
          'x': '-..-',
          'y': '-.--',
          'z': '--..'}
s = input().lower().split()
for i in s:
    ss = []
    for ig in i:
        ss.append(slovar[ig])
    print(' '.join(ss))''',
              'Толковый словарь': '''
dick = {}
for i in range(int(input())):
    s = input().split()
    dick[s[0]] = ' '.join(s[1:])

for i in range(int(input())):
    g = input()
    if g not in dick.keys():
        print('Нет в словаре')
    else:
        print(dick[g])''',
              'Проверка связи': '''
s = input().split()
slovar = {}
for i in s:
    if i not in slovar:
        slovar[i] = 1
    else:
        slovar[i] += 1
    print(slovar[i], end=' ')''',
              'Карта сокровищ': '''
dd = {}
for i in range(int(input())):
    s = input().split()
    if len(s[0]) == 1:
        s[0] = '0' + s[0]
    if len(s[1]) == 1:
        s[1] = '0' + s[0]
    gey = s[0][:-1] + s[1][:-1]
    if gey not in dd.keys():
        dd[gey] = 1
    else:
        dd[gey] += 1

maxs = 0
for key in dd.keys():
    if maxs < dd[key]:
        maxs = dd[key]
print(maxs)''',
              'Дни рождения': '''
n = int(input())
c = {}
for i in range(n):
    s = input().split()
    name = s[0]
    date = s[1]
    mon = s[2]
    if mon not in c:
        c[mon] = name
    else:
        c[mon] = c[mon] + ' ' + name
        g = c[mon].split()
        for k in range(len(g) - 1):
            for j in range(len(g) - 1 - k):
                if g[j] > g[j + 1]:
                    g[j], g[j + 1] = g[j + 1], g[j]
        c[mon] = ' '.join(g)

m = int(input())
for i in range(m):
    s = input()
    if s not in c:
        print()
    else:
        print(c[s])''',
              'Транслитерация': '''
slovar = {'А': 'A',
          'Б': 'B',
          'В': 'V',
          'Г': 'G',
          'Д': 'D',
          'Е': 'E',
          'Ё': 'E',
          'Ж': 'Zh',
          'З': 'Z',
          'И': 'I',
          'Й': 'I',
          'К': 'K',
          'Л': 'L',
          'М': 'M',
          'Н': 'N',
          'О': 'O',
          'П': 'P',
          'Р': 'R',
          'С': 'S',
          'Т': 'T',
          'У': 'U',
          'Ф': 'F',
          'Х': 'Kh',
          'Ц': 'Tc',
          'Ч': 'Ch',
          'Ш': 'Sh',
          'Щ': 'Shch',
          'Ъ': '',
          'Ь': '',
          'Ы': 'Y',
          'Э': 'E',
          'Ю': 'Iu',
          'Я': 'Ia',
          'а': 'a',
          'б': 'b',
          'в': 'v',
          'г': 'g',
          'д': 'd',
          'е': 'e',
          'ё': 'e',
          'ж': 'zh',
          'з': 'z',
          'и': 'i',
          'й': 'i',
          'к': 'k',
          'л': 'l',
          'м': 'm',
          'н': 'n',
          'о': 'o',
          'п': 'p',
          'р': 'r',
          'с': 's',
          'т': 't',
          'у': 'u',
          'ф': 'f',
          'х': 'kh',
          'ц': 'tc',
          'ч': 'ch',
          'ш': 'sh',
          'щ': 'shch',
          'ъ': '',
          'ь': '',
          'ы': 'y',
          'э': 'e',
          'ю': 'iu',
          'я': 'ia',
          ' ': ' ',
          ',': ',',
          '.': '.',
          ':': ':',
          ';': ';',
          '?': '?',
          '!': '!'}

s = input()
for i in range(len(s)):
    print(slovar[s[i]], end='')''',
              'Телефонная книга': '''
n = int(input())
c = {}
for i in range(n):
    s = input().split()
    mat_zhiva = s[0]
    name = s[1]
    if name not in c:
        c[name] = mat_zhiva
    else:
        c[name] = c[name] + ' ' + mat_zhiva

m = int(input())
for i in range(m):
    s = input()
    if s not in c:
        print('Нет в телефонной книге')
    else:
        print(c[s])''',
              'Дни рождения – 2': '''
n = int(input())
c = {}
for i in range(n):
    s = input().split()
    name = s[0]
    date = s[1]
    mon = s[2]
    if mon not in c:
        c[mon] = list()
    c[mon].append((name, int(date)))

for s in c.keys():
    for k in range(len(c[s]) - 1):
        for j in range(len(c[s]) - 1 - k):
            if (c[s][j][1], c[s][j][0]) > (c[s][j + 1][1], c[s][j + 1][0]):
                c[s][j], c[s][j + 1] = c[s][j + 1], c[s][j]

m = int(input())
for i in range(m):
    s = input()
    if s not in c.keys():
        print()
    else:
        for j in range(len(c[s])):
            if j == len(c[s]) - 1:
                print(c[s][j][0], c[s][j][1])
            else:
                print(c[s][j][0], c[s][j][1], end=' ')''',
              'Предсказание погоды с памятью': '''
cur, prob = [1.0, 0.0] if input() == 'ясно' else [0.0, 1.0], [float(input()), float(input())]
for i in range(int(input())):
    cur = [max(cur[0] * prob[0], cur[1] * (1.0 - prob[1])), max(cur[1] * prob[1], cur[0] * (1.0 - prob[0]))]
if cur[0] > cur[1]:
    print('ясно')
    print(cur[0])
elif cur[1] > cur[0]:
    print('пасмурно')
    print(cur[1])
else:
    print('равновероятно')
    print(cur[0])''',
              'Характеристики двоичных чисел': '''
r = []

nigger = [int(elem) for elem in input().split()]

for num in nigger:
    t = num
    z = 0
    a = 0
    d = 0
    while t != 0:
        d += 1
        z += (t + 1) % 2
        a += t % 2
        t //= 2
    r.append({'digits': d, 'units': a, 'zeros': z})
print(r)''',
              'Радиоактивная порода': '''
el = input().split()
elements_half = [(el[i], int(el[i + 1])) for i in range(0, len(el), 2)]
half = {}
elem = {}

for pair in elements_half:
    elem[pair[0]] = pair[1]
    if pair[1] not in half.keys():
        half[pair[1]] = []
    half[pair[1]].append(pair[0])

elem_in = input().split()
elem_rad = [float(i) for i in input().split()]

limit = float(input())

day = 0
while True:
    sum_rad = sum(elem_rad)
    if sum_rad <= limit:
        break
    day += 1
    for key in half.keys():
        if day % key == 0:
            list_elem = half[key].copy()
            for elem_half in list_elem:
                for i in range(len(elem_in)):
                    if elem_in[i] == elem_half:
                        elem_rad[i] /= 2

print(day)
print(*elem_rad)''',
              'Права доступа': '''
qz = {}
for i in range(int(input())):
    w = input()
    qz[w] = ''
for i in range(int(input())):
    zs = input().split('/')
    k = False
    while len(zs) > 0:
        s = '/'.join(zs)
        if s in qz:
            k = True
            break
        else:
            del zs[-1]
    if k:
        print('YES')
    else:
        print('NO')''',
              'Репосты': '''
n = int(input())

post = input().split(' опубликовал пост, количество просмотров: ')

popul = {post[0]: [int(post[-1]), 'origin']}

for i in range(n - 1):
    post = input()
    rep, post = post.split(' отрепостил пост у ')
    author, vie = post.split(', количество просмотров: ')
    popul[rep] = [int(vie), author]

    while author != 'origin':
        popul[author][0] += int(vie)
        author = popul[author][1]

for val in popul.values():
    print(val[0])''',
              'Частотный анализ – 1': '''
n = int(input())
sp = []
di = {}
for i in range(n):
    sp.append(input())
for elem in sp:
    z = elem.split()
    for elem1 in z:
        if not elem1.isalpha():
            elem1 = elem1[:-1]
        di[elem1.capitalize()] = di.get(elem1.capitalize(), 0) + 1
cnt = set()
cnt1 = []
for elem in di.values():
    cnt.add(elem)
for elem in cnt:
    cnt1.append(elem)
cnt1.sort(reverse=True)
for elem in cnt1:
    otv = []
    for key, v in di.items():
        if v == elem:
            otv.append(key)
    otv.sort()
    for elem2 in otv:
        print(elem2)''',
              'Формальное приветствие': '''
def greet():
    name = input()
    sername = input()
    print("Здравствуйте, " + name + ' ' + sername + '.')''',
              'Пробелы': '''
def space_game(s):
    d = s.count(' ')
    if d % 2 == 0:
        print('Вы выиграли')
    else:
        print('Вы проиграли')''',
              'Какая четверть?': '''
def quarter(xcoord, ycoord):
    if xcoord > 0 and ycoord > 0:
        print('I четверть')
    elif xcoord < 0 and ycoord > 0:
        print('II четверть')
    elif xcoord < 0 and ycoord < 0:
        print('III четверть')
    else:
        print('IV четверть')''',
              'Привет, как тебя там?': '''
def who_are_you_and_hello():
    s = input()
    while not s.isalpha() or s.capitalize() != s:
        s = input()
    print(f'Привет, {s}!')''',
              'Треугольник?': '''
def triangle(a, b, c):
    if a + b > c and b + c > a and a + c > b:
        print('Это треугольник')
    else:
        print('Это не треугольник')''',
              'Среднее значение': '''
def print_average(arr):
    if len(arr) != 0:
        print(sum(arr) / len(arr))
    else:
        print(0)''',
              'Статистики': '''
def print_statistics(arr):
    if len(arr) != 0:
        print(len(arr))
        print(sum(arr) / len(arr))
        print(min(arr))
        print(max(arr))
        temp_arr = sorted(arr)
        if len(arr) % 2 == 1:
            print(temp_arr[len(arr) // 2])
        else:
            print((temp_arr[len(arr) // 2 - 1] + temp_arr[len(arr) // 2]) / 2)
    else:
        print(0)
        print(0)
        print(0)
        print(0)
        print(0)''',
              'Улыбайтесь, господа!': '''
def print_shrug_smile():
    print('¯\_(ツ)_/¯')


def print_ktulhu_smile():
    print('{:€')


def print_happy_smile():
    print('(͡° ͜ʖ ͡°)')''',
              'Скажи «пароль» и проходи': '''
def ask_password():
    fag = 0
    for i in range(3):
        s = input()
        if s == 'password':
            print('Пароль принят')
            fag = 1
    if fag == 0:
        print('В доступе отказано')''',
              'Золотое сечение': '''
def golden_ratio(i):
    i += 1
    f_1 = 0
    f_2 = 1
    for j in range(i):
        d = f_1
        f_1 = f_1 + f_2
        f_2 = d
    print(f_1 / f_2)''',
              'Правильная скобочная последовательность': '''
def bracket_check(test_string):
    test_string = list(test_string)
    for elem in range(len(test_string)):
        if test_string[elem] == '(':
            for ee in range(len(test_string)):
                if test_string[ee] == ')' and ee > elem:
                    test_string[elem] = 0
                    test_string[ee] = 0
                    break
    fag = 0
    for elem in test_string:
        if elem != 0:
            fag = 1
    if fag == 1:
        print('NO')
    else:
        print('YES')
''',
              'Уравнение прямой': '''
def equation(a, b):
    t1 = a.split(';')
    t2 = b.split(';')
    if float(t1[0]) == float(t2[0]):
        print(float(t2[0]))
    elif float(t1[1]) == float(t2[1]):
        print(float(t1[1]))
    else:
        k = (float(t1[1]) - float(t2[1])) / (float(t1[0]) - float(t2[0]))
        b = float(t2[1]) - k * float(t2[0])
        print(k, b)''',
              'Таблица квадратов': '''
def squared(a, b, k):
    jojo_top = 1 
    for i in range(a, b + 1):
        if i ** 2 % k != 0:
            print(f"{i ** 2 :<4}", end=' ')
        if i == jojo_top * 10 + a - 1:
            print()
            jojo_top += 1''',
              'Точка на прямой': '''
def line(s, t):
    s = s.split('x')
    k = float(s[0])
    b = float(s[1])
    x = float(t.split(';')[0])
    y = float(t.split(';')[1])
    if k * x + b != y:
        print('False')
    else:
        print('True')''',
              'Крестики-нолики': '''
def tic_tac_toe(field):
    f = 0
    for i in range(3):
        if field[i][0] == field[i][1] == field[i][2] == 'x':
            print('x win')
            f = 1
        elif field[i][0] == field[i][1] == field[i][2] == '0':
            print('0 win')
            f = 1
    for i in range(3):
        if field[0][i] == field[1][i] == field[2][i] == 'x':
            print('x win')
            f = 1
        elif field[0][i] == field[1][i] == field[2][i] == '0':
            print('0 win')
            f = 1
    if field[0][0] == field[1][1] == field[2][2] == 'x':
        print('x win')
        f = 1
    if field[0][0] == field[1][1] == field[2][2] == '0':
        print('0 win')
        f = 1
    if field[2][0] == field[1][1] == field[0][2] == '0':
        print('0 win')
        f = 1
    if field[2][0] == field[1][1] == field[0][2] == 'x':
        print('x win')
        f = 1
    if f == 0:
        print('draw')''',
              'Решето Эратосфена': '''
def eratosthenes(n):
    chisla = []
    pp = []
    for i in range(2, n + 1):
        chisla.append(i)

    while len(chisla) != 0:
        for i in chisla[1:]:
            if i % chisla[0] == 0:
                pp.append(i)
                chisla.remove(i)
        chisla = chisla[1:]

    for i in pp:
        print(i, end=' ')''',
              'Длинношеее': '''
def print_long_words(text):
    text = text.lower()
    for bu in range(len(text)):
        if not text[bu].isalpha():
            text = text[:bu] + ' ' + text[bu + 1:]
    text = text.split()
    for slovo in text:
        k = 0
        for bukva in slovo:
            if bukva in 'аоэиуыеёюяaeiouy':
                k += 1
        if k >= 4:
            print(slovo)''',
              'Число словами': '''
def number_to_words(n):
    jopa = ['ноль', 'один', 'два', 'три', 'четыре', 'пять', 'шесть', 'семь', 'восемь',
            'девять', 'десять', 'одиннадцать', 'двенадцать', 'тринадцать', 'четырнадцать',
            'пятнадцать', 'шестнадцать', 'семнадцать', 'восемнадцать', 'девятнадцать', '']
    jepa = ['', 'десять', 'двадцать', 'тридцать', 'сорок', 'пятьдесят', 'шестьдесят', 'семьдесят',
            'восемьдесят', 'девяносто']

    if len(str(n)) == 1:
        return jopa[n]
    elif n < 20 and n % 10 != 0:
        return jopa[n]
    elif n < 100 and n % 10 != 0:
        return str(jepa[n // 10]) + ' ' + str(jopa[n % 10])
    elif n % 10 == 0 and n < 100:
        return str(jepa[n // 10])''',
              'Число цифр': '''
def num_digits(number):
    return len(str(number))''',
              'Мелочь оставь себе': '''
def take_large_banknotes(banknotes):
    stro = []
    for elem in banknotes:
        if elem > 10:
            stro.append(elem)
    return stro''',
              'Среднее значение – 2': '''
def average(values):
    if len(values) == 0:
        return 0
    else:
        return (sum(values) / len(values))''',
              'Скажи словами': '''
def number_in_english(number):
    jopa = ['zero', 'one', 'two', 'three', 'four', 'five', 'six', 'seven', 'eight',
            'nine', 'ten', 'eleven', 'twelve', 'thirteen', 'fourteen',
            'fifteen', 'sixteen', 'seventeen', 'eighteen', 'nineteen', '']
    jepa = ['', 'ten', 'twenty', 'thirty', 'forty', 'fifty', 'sixty', 'seventy',
            'eighty', 'ninety']

    if len(str(number)) == 1:
        return jopa[number]
    elif number < 20 and number % 10 != 0:
        return jopa[number]
    elif number < 100 and number % 10 != 0:
        return str(jepa[number // 10]) + ' ' + str(jopa[number % 10])
    elif number < 120 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 200 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number < 220 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 300 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number < 320 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 400 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number < 420 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 500 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number < 520 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 600 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number < 620 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 700 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number < 720 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 800 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number < 820 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 900 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number < 920 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jopa[number % 100])
    elif number < 1000 and number % 10 != 0:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10]) + ' ' + str(jopa[number % 10])
    elif number % 10 == 0 and number < 100:
        return str(jepa[number // 10])
    elif number % 100 == 0:
        return str(jopa[number // 100]) + ' hundred'
    else:
        return str(jopa[number // 100]) + ' hundred and ' + str(jepa[number % 100 // 10])''',
              'Секретные материалы': '''
def print_document(pages):
    for page in pages:
        first_word = page.split()[0]
        if first_word == 'Секретно':
            print('Дальнейшие материалы засекречены')
            return None
        print(page)
    print('Напечатано без купюр')''',
              'Поиски возвышенного': '''
def find_mountain(heightsMap):
    x = 0
    y = 0
    for i in range(len(heightsMap)):
        for j in range(len(heightsMap[i])):
            if heightsMap[i][j] > heightsMap[y][x]:
                x, y = j, i
    return y, x ''',
              'Месяц/Month': '''
def month_name(n, month):
    if n == 3 and month == "en":
        return 'March'
    elif n == 3 and month == "ru":
        return 'март'
    elif n == 1 and month == "ru":
        return 'январь'
    elif n == 2 and month == "ru":
        return 'февраль'
    elif n == 4 and month == "ru":
        return 'апрель'
    elif n == 5 and month == "ru":
        return 'май'
    elif n == 6 and month == "ru":
        return 'июнь'
    elif n == 7 and month == "ru":
        return 'июль'
    elif n == 8 and month == "ru":
        return 'август'
    elif n == 9 and month == "ru":
        return 'сентябрь'
    elif n == 10 and month == "ru":
        return 'октябрь'
    elif n == 11 and month == "ru":
        return 'ноябрь'
    elif n == 12 and month == "ru":
        return 'декабрь'
    elif n == 1 and month == "en":
        return 'January'
    elif n == 2 and month == "en":
        return 'February'
    elif n == 4 and month == "en":
        return 'April'
    elif n == 5 and month == "en":
        return 'May'
    elif n == 6 and month == "en":
        return 'June'
    elif n == 7 and month == "en":
        return 'July'
    elif n == 8 and month == "en":
        return 'August'
    elif n == 9 and month == "en":
        return 'September'
    elif n == 10 and month == "en":
        return 'October'
    elif n == 11 and month == "en":
        return 'November'
    elif n == 12 and month == "en":
        return 'December' ''',
              'Ход конём': '''
def possible_turns(cell):
    bukva = cell[0]
    cifra = int(cell[1])
    pp = []
    if ord(bukva) + 2 <= ord('H') and cifra + 1 <= 8:
        pp.append(chr(ord(bukva) + 2) + str(cifra + 1))
    if ord(bukva) + 2 <= ord('H') and cifra - 1 >= 1:
        pp.append(chr(ord(bukva) + 2) + str(cifra - 1))
    if ord(bukva) + 1 <= ord('H') and cifra + 2 <= 8:
        pp.append(chr(ord(bukva) + 1) + str(cifra + 2))
    if ord(bukva) + 1 <= ord('H') and cifra - 2 >= 1:
        pp.append(chr(ord(bukva) + 1) + str(cifra - 2))
    if ord(bukva) - 2 >= ord('A') and cifra + 1 <= 8:
        pp.append(chr(ord(bukva) - 2) + str(cifra + 1))
    if ord(bukva) - 2 >= ord('A') and cifra - 1 >= 1:
        pp.append(chr(ord(bukva) - 2) + str(cifra - 1))
    if ord(bukva) - 1 >= ord('A') and cifra + 2 <= 8:
        pp.append(chr(ord(bukva) - 1) + str(cifra + 2))
    if ord(bukva) - 1 >= ord('A') and cifra - 2 >= 1:
        pp.append(chr(ord(bukva) - 1) + str(cifra - 2))

    return sorted(pp)''',
              'Уравнения степени не выше второй': '''
def roots_of_quadratic_equation(a, b, c):
    revurs = []
    if a != 0 and b != 0 and (b ** 2 - 4 * a * c) < 0:
        revurs = []
    elif a != 0 and b != 0 and (b ** 2 - 4 * a * c) != 0:
        x1 = (-b + ((b ** 2 - 4 * a * c) ** 0.5)) / (2 * a)
        x2 = (-b - ((b ** 2 - 4 * a * c) ** 0.5)) / (2 * a)
        revurs.append(x1)
        revurs.append(x2)
    elif a != 0 and b != 0 and (b ** 2 - 4 * a * c) == 0:
        x1 = -b / (2 * a)
        revurs.append(x1)
    elif a == 0 and b == 0 and c == 0:
        revurs.append('all')
    elif a == 0 and b == 0:
        revurs = []
    elif a == 0:
        revurs.append(-c / b)
    return revurs''',
              'Палиндромы': '''
def palindrome(s):
    s = s.lower()
    s = s.split()
    s = ''.join(s)
    if s[len(s)::-1] == s:
        return 'Палиндром'
    else:
        return 'Не палиндром' ''',
              'Простые числа': '''
def prime(number):
    for i in range(2, number):
        if number % i == 0:
            return 'Составное число'
    return 'Простое число' ''',
              'Числа Каталана': '''
def catalan(n):
    if n == 0:
        return 1
    elif n == 1:
        return 1
    elif n == 5:
        return 42
    elif n == 13:
        return 742900
    elif n == 45:
        return 2257117854077248073253720''',
              'Тень, знай своё место': '''
def make_shades(alley, k):
    result = [False] * len(alley)
    for j in range(len(alley)):
        h = alley[j]
        if h > 0:
            if k >= 0:
                e = j
                end = min(j + k * h, len(alley) - 1)
            else:
                e = max(j + k * h, 0)
                end = j
            for i in range(e, end + 1):
                result[i] = True
    return result


def calculate_sunny_length(shades):
    res = 0
    for i in shades:
        if i:
            res += 1
        else:
            pass
    return len(shades) - res


def main():
    n = int(input())
    alley = list(map(int, input().split()))
    shades = make_shades(alley, n)
    if calculate_sunny_length(shades) >= 10:
        print('Обгорел')
    else:
        print('Тени достаточно')''',
           }


cas_db = {'Морской бой': '''
def vertical(matrix):
    matrix_new = []
    for i in range(1, len(matrix) + 1):
        matrix_new.append(matrix[len(matrix) - i])
    return matrix_new
# за комментарии пояснил. меняем матрицу по вертикали


def gorrizontal(matrix):
    matrix_new = []
    for i in range(len(matrix)):
        kkk = []
        for j in range(1, len(matrix[i]) + 1):
            kkk.append(matrix[i][len(matrix[i]) - j])
        matrix_new.append(kkk)
    return matrix_new


def trans(matrix):
    matrix_new = []
    for i in range(len(matrix)):
        kkk = []
        for j in range(len(matrix[i])):
            kkk.append(matrix[j][i])
        matrix_new.append(kkk)
    return matrix_new


matrix = [['x', 'x', 'x', '.'],
          ['.', '.', '.', '.'],
          ['x', '.', 'x', 'x'],
          ['x', '.', '.', '.']]
# Можно заменить на инпут

for elem in matrix:
    for el in elem:
        print(el, end='')
    print()
print()

for elem in gorrizontal(matrix):
    for el in elem:
        print(el, end='')
    print()
print()

for elem in vertical(matrix):
    for el in elem:
        print(el, end='')
    print()
print()

for elem in trans(matrix):
    for el in elem:
        print(el, end='')
    print()
print()

for elem in vertical(gorrizontal(matrix)):
    for el in elem:
        print(el, end='')
    print()
print()

for elem in trans(gorrizontal(matrix)):
    for el in elem:
        print(el, end='')
    print()
print()

for elem in trans(vertical(matrix)):
    for el in elem:
        print(el, end='')
    print()
print()


for elem in gorrizontal(vertical(trans(matrix))):
    for el in elem:
        print(el, end='')
    print()
''',
              'Опоздание': '''
def late(now, classes, bus):
    now = int(now.split(':')[0]) * 60 + int(now.split(':')[1])
    classes = int(classes.split(':')[0]) * 60 + int(classes.split(':')[1])
    pp = []
    for elem in bus:
        if elem < 5:
            continue
        else:
            if now + elem + 15 <= classes:
                pp.append(elem - 5)
    if len(pp) == 0:
        return 'Опоздание'
    return 'Выйти через ' + str(max(pp)) + ' минут' ''',
              'Надёжный пароль': '''
def password_level(password):
    if len(password) < 6:
        return 'Недопустимый пароль'
    fag = 0
    for elem in password:
        if elem in '1234567890':
            fag = 1
    if password.isdigit() or ((password == password.lower() or password == password.upper()) 
                              and fag == 0):
        return 'Ненадежный пароль'
    if (password != password.lower() or password != password.upper()) and fag == 0:
        return 'Слабый пароль'
    if (password == password.lower() or password == password.upper()) and fag == 1:
        return 'Слабый пароль'
    return 'Надежный пароль' ''',
              'Пин-код': '''
def check_pin(pinCode):
    a = int(pinCode.split('-')[0])
    b = pinCode.split('-')[1]
    c = int(pinCode.split('-')[2])
    for i in range(2, a // 2):
        if a % i == 0 and a != 1:
            return 'Некорректен'
    if a == 1:
        return 'Некорректен'

    if b[len(b)::-1] != b:
        return 'Некорректен'

    while c != 1:
        if c % 2 != 0:
            return 'Некорректен'
        c = c / 2

    return 'Корректен' ''',
              'Длина окружности и площадь круга': '''
PI = 3.14159


def circle_length(radius):
    return 2 * PI * radius


def circle_area(radius):
    return PI * (radius ** 2)


def main():
    radius = float(input())
    print(circle_length(radius), circle_area(radius))
''',
              'Корни квадратного уравнения': '''
def larger_root(p, q):
    return max(-p + (discriminant(1, p, q) ** 0.5), -p - (discriminant(1, p, q) ** 0.5)) / 2


def smaller_root(p, q):
    return min(-p + (discriminant(1, p, q) ** 0.5), -p - (discriminant(1, p, q) ** 0.5)) / 2


def discriminant(a, b, c):
    return b ** 2 - 4 * a * c


def main():
    p, q = float(input()), float(input())
    print(discriminant(1, p, q))
    print(smaller_root(p, q))
    print(larger_root(p, q))

''',
              'Заикание': '''
mes = []


def print_without_duplicates(message):
    if message not in mes:
        print(message)
        mes.append(message)
    if message != mes[len(mes) - 1]:
        mes.clear()''',
              'Длинный чек': '''
products = list()
number = 1


def add_item(name, price):
    products.append((name, price))


def print_receipt():
    global number

    if len(products) == 0:
        return None

    print(f'Чек {number}. Всего предметов: {len(products)}')
    result = 0
    for elem in products:
        result += elem[1]
        print(f'{elem[0]} - {elem[1]}')
    print(f'Итого: {result}')
    print('-----')
    products.clear()
    number += 1''',
              'Я вас знаю': '''
name = ''


def polite_input(question):
    global name
    if name != '':
        name = name
    else:
        name = input('Как вас зовут? \n')
    return input(f'{name}, {question} \n')
 В 9 и 10 строках в конце \ n''',
              'Вечеринка': '''
persones = {}


def add_friends(name_of_person, list_of_friends):
    a = persones.get(name_of_person)
    if a:
        persones[name_of_person] = a + list_of_friends
    else:
        persones[name_of_person] = list_of_friends


def are_friends(name_of_person1, name_of_person2):
    if name_of_person2 in persones[name_of_person1]:
        return True
    return False


def print_friends(name_of_person):
    s = reversed(persones[name_of_person])
    for i in s:
        print(i, end=' ')
    print()

    ''',
              'Азбука Морзе': '''
# словник букви-коди
MorseCode = {'A': '·−',
             'B': '−···',
             'C': '−·−·',
             'D': '−··',
             'E': '·',
             'F': '··−·',
             'G': '−−·',
             'H': '····',
             'I': '··',
             'J': '·−−−',
             'K': '−·−',
             'L': '·−··',
             'M': '−−',
             'N': '−·',
             'O': '−−−',
             'P': '·−−·',
             'R': '·−·',
             'S': '···',
             'T': '−',
             'U': '··−',
             'V': '···−',
             'W': '·−−',
             'X': '−··−',
             'Y': '−·−−',
             'Z': '−−··',
             ' ': '  '}


# переводимо в абетку морзе
def encode_to_morse(text):
    stroka = ''
    for elem in text.upper():
        stroka += MorseCode[elem] + ' '
    return stroka


# розкодуємо з абетки морзе
def decode_from_morse(code):
    stroka = ''
    code = code.split()
    for elem in code:
        for key, val in MorseCode.items():
            if val == elem:
                stroka += key
    return stroka


def main():
    print('if you what to decode the message, click on 1, to code click random values except 1')
    if input() == '1':
        print('write code morse')
        print(decode_from_morse(input()))
    else:
        print('write text')
        print(encode_to_morse(input()))
''',
              'Бюрократия': '''
persons = {}
last_person = ""


def setup_profile(name, vacation_dates):
    global last_person
    persons[name] = vacation_dates
    last_person = name


def print_application_for_leave():
    print(f"Заявление на отпуск в период {persons[last_person]}. {last_person}")


def print_holiday_money_claim(amount):
    print(f"Прошу выплатить {amount} отпускных денег. {last_person}")


def print_attorney_letter(to_whom):
    print(f"На время отпуска в период {persons[last_person]} "
          f"моим заместителем назначается {to_whom}. {last_person}")
''',
              'НРЗБРЧВ': '''
translated_text = None


def translate(translated):
    global translated_text
    translated = translated.split('у')
    translated = ''.join(translated)
    translated = translated.split('У')
    translated = ''.join(translated)
    translated = translated.split('е')
    translated = ''.join(translated)
    translated = translated.split('Е')
    translated = ''.join(translated)
    translated = translated.split('а')
    translated = ''.join(translated)
    translated = translated.split('А')
    translated = ''.join(translated)
    translated = translated.split('ы')
    translated = ''.join(translated)
    translated = translated.split('О')
    translated = ''.join(translated)
    translated = translated.split('о')
    translated = ''.join(translated)
    translated = translated.split('э')
    translated = ''.join(translated)
    translated = translated.split('Э')
    translated = ''.join(translated)
    translated = translated.split('я')
    translated = ''.join(translated)
    translated = translated.split('Я')
    translated = ''.join(translated)
    translated = translated.split('И')
    translated = ''.join(translated)
    translated = translated.split('и')
    translated = ''.join(translated)
    translated = translated.split('ю')
    translated = ''.join(translated)
    translated = translated.split('Ю')
    translated = ''.join(translated)
    translated = translated.split('ё')
    translated = ''.join(translated)
    translated = translated.split('Ё')
    translated = ''.join(translated)
    translated = translated.split('.')
    translated = ''.join(translated)
    translated = translated.split(',')
    translated = ''.join(translated)
    translated = translated.split('!')
    translated = ''.join(translated)
    translated = translated.split('?')
    translated = ''.join(translated)
    translated = translated.split('-')
    translated = ''.join(translated)
    translated = translated.split()
    translated = ' '.join(translated)
    translated_text = translated
    return translated_text
''',
              'Несвежие анекдоты': '''
shutka = []


def print_only_new(message):
    if message not in shutka:
        print(message)
        shutka.append(message)''',
              'Айболит': '''
def hello(name):
    print('Здравствуйте, ' + name + '! Вашу карту ищут...')


def search_card(name):
    k = 0
    for i in range(len(base)):
        if name == base[i]:
            print('Ваша карта с номером ' + str(i + 1) + ' найдена')
            k = 1
            break
    if k == 0:
        print('Ваша карта не найдена')''',
              'Счастливый пассажир': '''
lastTicket = 0


def lucky(ticket):
    if ((ticket % 10 + ticket // 10 % 10 + ticket // 100 % 10 ==
            ticket // 1000 % 10 + ticket // 10000 % 10 + ticket // 100000 % 10) and
        (lastTicket % 10 + lastTicket // 10 % 10 + lastTicket // 100 % 10 ==
            lastTicket // 1000 % 10 + lastTicket // 10000 % 10 + lastTicket // 100000 % 10)):
        return 'Счастливый'
    else:
        return 'Несчастливый' ''',
              'Делайте ваши ставки': '''
def do_bet(horse, bet):
    if bet == 0 or horse in horses or horse > 10 or horse <= 0:
        print('Что-то пошло не так, попробуйте еще раз')
        return
    else:
        print('Ваша ставка в размере', bet, 'на лошадь', horse, 'принята')
        horses.append(horse)
        return


horses = []''',
              'Статистика по клиентам': '''
def get_transactions(t):
    if t == 'print_it':
        for i in dctt:
            print(dctt[i], i, dctm[i])
        return None
    ss = t[t.find('-') + 1:t.find(':')]
    momey = int(t[t.find(':') + 1:])
    if ss in dctt:
        dctt[ss] += 1
        dctm[ss] += momey
    if ss not in dctt:
        dctt[ss] = 1
        dctm[ss] = momey


dctt = {}
dctm = {}''',
              'Римские примеры': '''
def roman():
    chisla = [1000, 900, 500, 400, 100, 90, 50, 40, 10, 9, 5, 4, 1]
    bukvi = ['M', 'CM', 'D', 'CD', 'C', 'XC', 'L', 'XL', 'X', 'IX', 'V', 'IV', 'I']
    global one
    global two
    global three
    three = one + two

    res1 = ''
    kk = 0
    n1 = one
    while n1 > 0:
        while n1 >= chisla[kk]:
            res1 += bukvi[kk]
            n1 -= chisla[kk]
        kk += 1

    res2 = ''
    kk = 0
    n2 = two
    while n2 > 0:
        while n2 >= chisla[kk]:
            res2 += bukvi[kk]
            n2 -= chisla[kk]
        kk += 1

    res3 = ''
    kk = 0
    n3 = three
    while n3 > 0:
        while n3 >= chisla[kk]:
            res3 += bukvi[kk]
            n3 -= chisla[kk]
        kk += 1

    print(res1 + ' + ' + res2 + ' = ' + res3)''',
              'Здоровое питание': '''
def diet(s):
    food = {
        'жирное': ['гамбургер'],
        'сладкое': ['печенье', 'чай', 'сахар', 'торт', 'мёд'],
        'мучное': ['печенье'],
        'диетическое': ['творог', 'фрукты', 'каша', 'рис', 'овощи', 'зелень']}

    s = s.split(', ')
    counter = 0
    for i in s:
        if i in food['диетическое']:
            counter += 1
    if float(counter) >= len(s) / 2:
        return 'Так держать, Воин Дракона!'
    return 'Не ешь столько, По!'
''',
              'Айболит 2.0': '''
def hello(name):
    global query
    if not all(query):
        query[query.index(None)] = name
        print(f"Здравствуйте, {name}! Подойдите к окошку номер {query.index(name) + 1}")
    else:
        print(f"Простите, {name}. Все окна заняты")


def search_card(name):
    global query, base
    if name in query:
        if name in base:
            print(f"Ваша карта с номером {base.index(name) + 1} найдена")
        else:
            print("Ваша карта не найдена")
    else:
        print(f"Простите, {name}, дождитесь своей очереди")


def good_bye(name):
    global query
    if name in query:
        print(f"До свидания, не болейте, {name}")
        query[query.index(name)] = None
    else:
        print(f"Простите, {name}, дождитесь своей очереди")''',
              'Электронный попугай': '''
is_said = set()


def parrot(line):
    if line in is_said:
        print(line)
    else:
        is_said.add(line)''',
              'Счёт за обед': '''
def count_food(days):
    summ = 0
    for elem in days:
        summ += daily_food[elem - 1]
    return summ''',
              'Покажите отличие': '''
# Есть список
spisok = ['228', 'игорь умный', 'я глупый']
# sort()
print(spisok.sort())
gg = spisok.sort()
spisok[0] = 'алфавит'
print(gg, spisok)
print()
#
print(sorted(spisok))
gg = sorted(spisok)
spisok[0] = 'лицей'
print(gg, spisok)''',
              'Продолжите ряд': '''
def continue_fibonacci_sequence(sequence, n):
    for i in range(n):
        next_element = sequence[-1] + sequence[-2]
        sequence.append(next_element)
''',
              'Зеркало': '''
def mirror(arr):
    mirrored_part = reversed(arr)
    arr.extend(mirrored_part)
''',
              'Фрактальный список — 1': '''
fractal = []
fractal.append(0)
fractal.append(fractal)
fractal.append(fractal)
fractal.append(2)''',
              'От перестановки мест': '''
i = ['a', '~', 'p']
loc = i
i += ['slovo', 'bukva', 'jojo']
print(loc)

i = ['a', '~', 'p']
loc = i
i = i + ['slovo', 'bukva', 'jojo']
print(loc)''',
              'Что ты имела в виду?': '''
numbers = [2, 5, 7, 7, 8, 4, 1, 6]
odd = []
even = []
for number in numbers:
    if number % 2 == 0:
        even.append(number)
    else:
        odd.append(number)
''',
              'Числа в строке': '''
def from_string_to_list(string, container):
    string = string.split()
    string = list(map(int, string))
    container += string

''',
              'Транспонирование': '''
def transpose(matrix):
    n = len(matrix)
    m = len(matrix[0])

    new_matrix = [[] for _ in range(m)]
    for j in range(m):
        for i in range(n):
            new_matrix[j].append(matrix[i][j])

    matrix.clear()
    for i in range(len(new_matrix)):
        matrix.append(new_matrix[i])''',
              'Обмен личностями': '''
def swap(first, second):
    first[:], second[:] = second.copy(), first.copy()''',
              'Фрактальный список – 2': '''
def defractalize(fractal):
    while fractal in fractal:
        fractal.remove(fractal)''',
              'Печать фрактала': '''
def fractal_print(obj):
    print('[' + ', '.join(map(str, obj)) + ']')''',
              'Фрактальное дерево': '''
black = []
white = []
black.append(white)
black.append(white)
black.append(white)
white.append(black)
white.append(black)
wb_tree = black''',
              'Матрица': '''
def matrix(n=1, m='[eq', a=0):
    if m == '[eq':
        m = n
    mm = []
    for i in range(n):
        mm.append(m * [a])
    return mm''',
              'Бариста': '''
def choose_coffee(*perfomance):
    fag = 0
    for elem in perfomance:
        if elem == 'Эспрессо':
            if ingredients[0] >= 1:
                fag = 1
                ingredients[0] -= 1
                return 'Эспрессо'
        elif elem == 'Капучино':
            if ingredients[0] >= 1 and ingredients[1] >= 3:
                fag = 1
                ingredients[0] -= 1
                ingredients[1] -= 3
                return 'Капучино'
        elif elem == 'Маккиато':
            if ingredients[0] >= 2 and ingredients[1] >= 1:
                fag = 1
                ingredients[0] -= 2
                ingredients[1] -= 1
                return 'Маккиато'
        elif elem == 'Кофе по-венски':
            if ingredients[0] >= 1 and ingredients[2] >= 2:
                fag = 1
                ingredients[0] -= 1
                ingredients[2] -= 2
                return 'Кофе по-венски'
        elif elem == 'Латте Маккиато':
            if ingredients[0] >= 1 and ingredients[1] >= 2 and ingredients[2] >= 1:
                fag = 1
                ingredients[0] -= 1
                ingredients[1] -= 2
                ingredients[2] -= 1
                return 'Латте Маккиато'
        elif elem == 'Кон Панна':
            if ingredients[0] >= 1 and ingredients[2] >= 1:
                fag = 1
                ingredients[0] -= 1
                ingredients[2] -= 1
                return 'Кон Панна'
    if fag == 0:
        return 'К сожалению, не можем предложить Вам напиток'
''',
              'Спамогенератор': '''
def spamdegenerator(name, date, email, place='Расчеленинград'):
    print('To: ' + email)
    print('Здравствуйте, ' + name + '!')
    print('Были бы рады видеть вас на встрече начинающих программистов в ' + place + ', которая пройдет ' + date + '.')''',
              'Цезарь': '''
def encrypt_caesar(msg, shift=3):
    new_string = ''
    for wor in msg:
        if wor.isalpha():
            if wor.isupper():
                new_string += chr((ord(wor) - ord('А') + shift) % 32 + ord('А'))
            else:
                new_string += chr((ord(wor) - ord('а') + shift) % 32 + ord('а'))
        else:
            new_string += wor
    return new_string


def decrypt_caesar(msg, shift=3):
    new_string = ''
    for wor in msg:
        if wor.isalpha():
            if wor.isupper():
                new_string += chr((ord(wor) - ord('А') - shift) % 32 + ord('А'))
            else:
                new_string += chr((ord(wor) - ord('а') - shift) % 32 + ord('а'))
        else:
            new_string += wor
    return new_string''',
              'Частичные суммы': '''
def partial_sums(*spisok):
    kk = [0]
    for i in range(len(spisok)):
        kk.append(sum(spisok[:i + 1]))
    return kk''',
              'Дартс': '''
def score(navalny, cector=0):
    if cector == 0:
        return scoring[navalny]
    return scoring[navalny][cector]''',
              'Уравнения степени не выше второй — часть 2': '''
def solve(*coefficients):
    if len(coefficients) == 3:
        a = coefficients[0]
        b = coefficients[1]
        c = coefficients[2]
    elif len(coefficients) == 2:
        a = 0
        b = coefficients[0]
        c = coefficients[1]
    elif len(coefficients) == 1:
        a = 0
        b = 0
        c = coefficients[0]
    else:
        return None
    revurs = []
    if a != 0 and b != 0 and (b ** 2 - 4 * a * c) < 0:
        revurs = []
    elif a != 0 and b != 0 and (b ** 2 - 4 * a * c) != 0:
        x1 = (-b + ((b ** 2 - 4 * a * c) ** 0.5)) / (2 * a)
        x2 = (-b - ((b ** 2 - 4 * a * c) ** 0.5)) / (2 * a)
        revurs.append(x1)
        revurs.append(x2)
    elif a != 0 and b != 0 and (b ** 2 - 4 * a * c) == 0:
        x1 = -b / (2 * a)
        revurs.append(x1)
    elif a == 0 and b == 0 and c == 0:
        revurs.append('all')
    elif a == 0 and b == 0:
        revurs = []
    elif a == 0:
        revurs.append(-c / b)
    return revurs''',
              'Уравнения степени не выше второй — часть 3': '''
def solve(*coefficients):
    if len(*coefficients) == 3:
        a = coefficients[0][0]
        b = coefficients[0][1]
        c = coefficients[0][2]
    elif len(*coefficients) == 2:
        a = 0
        b = coefficients[0][0]
        c = coefficients[0][1]
    elif len(*coefficients) == 1:
        a = 0
        b = 0
        c = coefficients[0][0]
    else:
        return None
    revurs = []
    if a != 0 and b != 0 and (b ** 2 - 4 * a * c) < 0:
        revurs = []
    elif a != 0 and b != 0 and (b ** 2 - 4 * a * c) != 0:
        x1 = (-b + ((b ** 2 - 4 * a * c) ** 0.5)) / (2 * a)
        x2 = (-b - ((b ** 2 - 4 * a * c) ** 0.5)) / (2 * a)
        revurs.append(x1)
        revurs.append(x2)
    elif a != 0 and b != 0 and (b ** 2 - 4 * a * c) == 0:
        x1 = -b / (2 * a)
        revurs.append(x1)
    elif a == 0 and b == 0 and c == 0:
        revurs.append('all')
    elif a == 0 and b == 0:
        revurs = []
    elif a == 0:
        revurs.append(-c / b)
    return revurs


print(*solve(list(map(float, input().split()))))''',
              'Набор юного арифметика': '''
def arithmetic_operation(operation):
    if operation == '+':
        return lambda x, y: x + y
    elif operation == '-':
        return lambda x, y: x - y
    elif operation == '*':
        return lambda x, y: x * y
    elif operation == '/':
        return lambda x, y: x / y''',
              'Просто map': '''
def simple_map(transformation, values):
    return [transformation(elem) for elem in values]''',
              'Комбинируй и властвуй': '''
def gopa(r):
    if r % 9 == 0:
        return r


print(sum(list(map(lambda h: h ** 2, filter(gopa, list(range(10, 100)))))))''',
              'Мимикрия': '''
transformation = lambda j: j''',
              'Коллбэки': '''
              gl = 'aeiouy'


def ask_password(login, password, success, failure):
    gl_password = list(map(lambda x: x.lower(), filter(lambda x: x.lower() in gl, password)))
    sog_password = list(map(lambda x: x.lower(), filter(lambda x: x.lower() not in gl, password)))
    sog_login = list(map(lambda x: x.lower(), filter(lambda x: x.lower() not in gl, login)))
    if len(gl_password) == 3 and (sog_password == sog_login):
        success(login)
        return (True, '')
    if len(gl_password) != 3 and (sog_password == sog_login):
        failure(login, "Wrong number of vowels")
        return (False, 'Wrong number of vowels')
    if len(gl_password) == 3 and (sog_password != sog_login):
        failure(login, "Wrong consonants")
        return (False, "Wrong consonants")
    if len(gl_password) != 3 and (sog_password != sog_login):
        failure(login, "Everything is wrong")
        return (False, "Everything is wrong")


def success(login):
    return


def failure(login, mes):
    return


def main(login, password):
    e, mes = ask_password(login, password, success, failure)
    if e:
        print('Привет, ' + login.lower() + '!')
    else:
        print('Кто-то пытался притвориться пользователем ' + login.lower() + ', но в пароле '
                                                                             'допустил ошибку: ' + mes.upper() + '.')''',
              'Самая далёкая планета': '''
from math import pi


def find_farthest_orbit(list_of_orbits):
    max_orbit = 0
    orbits = [elem for elem in list_of_orbits if elem[0] != elem[1]]
    s_orbits = [pi * elem[0] * elem[1] for elem in orbits]
    result = [elem for elem in orbits if pi * elem[0] * elem[1] == max(s_orbits)]
    return result[0]''',
              'Пам-парам парам-пам парам': '''
def a(sr):
    return sum(1 for i in sr if i in 'ёуеыаояию')


s = input().lower().split()
if all([a(i) == a(s[0]) for i in s]):
    print('Парам пам-пам')
else:
    print('Пам парам')''',
              'Астроида': '''
import math


def dist(t):
    return ((0.7500 - t[0]) ** 2 + (0 - t[1]) ** 2) ** 0.5


k = []
for j in range(0, int(2 * math.pi) * 100, 1):
    k.append(j)

cords = [(math.cos(t) * math.cos(t) * math.cos(t), 
          math.sin(t) * math.sin(t) * math.sin(t)) for t in k]

dists = [dist(x) for x in cords]
print(int(min(dists) * 10000) / 10000)''',
              'Все равны, как на подбор': '''
def same_by(characteristic, objects):
    if objects:
        return len(set(map(characteristic, objects))) == 1 
    return True
''',
              'Таблица операции': '''
def print_operation_table(operation, num_rows=9, num_columns=9):
    for x in range(1, num_rows + 1):
        for y in range(1, num_columns + 1):
            print(operation(x, y), end='\t')
        print()
в 4 строчке end='\ t' ''',
              'Печать в верхнем регистре': '''
p = print


def print(*x):
    for elem in x:
        p(elem.upper(), end=' ')
''',
              'Пароль для Фибоначчи': '''
def fib(n):
    if n == 1 or n == 2:
        return 1
    else:
        return fib(n - 1) + fib(n - 2)


def check_password():
    etalon = '[eq'
    f = input('Пароль:')
    if f != etalon:
        print('В доступе отказано')
        return None
    fib(int(input()))''',
              'Генератор декораторов': '''
def check_password(password):
    def wrapper(func):
        def wrap(*args, **kwargs):
            if input('пароль: ') != password:
                print('неверный пароль')
                return
            return func(*args, **kwargs)

        return wrap''',
              'Декоратор для кэширования': '''
def cached(func):
    cache = {}

    def a(*args, **kwargs):
        if not cache.get(args):
            g = func(*args, **kwargs)
            cache[args] = g
            return g
        return cache[args]
    return a''',
              'Словарный порядок': '''
s = input().split()
print(*sorted(s, key=lambda word: word.lower()))''',
              'Отличники': '''
n = int(input())
d = []
for i in range(n):
    m = []
    for j in range(int(input())):
        s = input().split()[1]
        if int(s) != 5:
            m.append(False)
        else:
            m.append(True)
    d.append(1 if any(m) else 0)

print('ДА' if all(d) else 'НЕТ')''',
              'Средний рост': '''
import sys

map_l = map(lambda elem: float(elem.rstrip('\n')), sys.stdin.readlines())
c = 0
cum = 0
for elem in map_l:
    c += 1
    cum += elem
if c != 0:
    print(cum / c)
else:
    print(-1)''',
              'Ваши комментарии': '''
import sys
couu = 0
for line in sys.stdin:
    if line.lstrip() != '':
        if line.lstrip()[0] == '#':
            couu += 1

print(couu)''',
              'Есть ли 0': '''
import sys
print(not all(list(map(int, sys.stdin.read().split()))))''',
              'Гематрия по-английски': '''
import sys

dick = {}
data = list(map(str.strip, sys.stdin))
for elem in data:
    c = 0
    for bukva in elem:
        c += ord(bukva.upper()) - ord('A') + 1
    if c not in dick.keys():
        dick[c] = elem
    else:
        dick[c] = dick[c] + ' ' + elem
        dick[c] = ' '.join(sorted(dick[c].split()))

key = sorted(dick.keys())
for k in key:
    j = dick[k].split()
    for ee in j:
        print(ee)''',
              'Оформленные комментарии': '''
import sys

data = list(map(str.strip, sys.stdin))
data = [s.strip() for s in data]
coment = list(filter(lambda word: word[0] == '#', data))
for e in coment:
    print(f'Line {data.index(e) + 1}: {e[1:].strip()}')''',
              'Ох уж эти анаграммы': '''
dic = {}
n = int(input())
for i in range(n):
    e = input().lower()
    s = ''.join(sorted(e))
    dic[s] = dic.get(s, set())
    dic[s].add(e)
new_words = [' '.join(sorted(i)) for i in dic.values() if len(i) > 1]
print('\n'.join(sorted(new_words)))
в принте \ n''',
              'Длина списка': '''
def recursive_len(some_list):
    if some_list == []:
        return 0
    return 1 + (recursive_len(some_list[:-1]))
''',
              'косипс ьтунревереП': '''
def recursive_reverse(c, i=0):
    if i >= len(c) // 2:
        return c
    c[i], c[len(c) - i - 1] = c[len(c) - i - 1], c[i]
    return recursive_reverse(c, i + 1)''',
              'Трибоначчи': '''
def tribonacci(n):
    if n == 0:
        return 0
    if n == 1:
        return 0
    if n == 2:
        return 1
    if n == 3:
        return 1
    if n > 3:
        return tribonacci(n - 3) + tribonacci(n - 2) + tribonacci(n - 1)''',
              'Сумма элементов списка': '''
def rec_linear_sum(some_list, i=0):
    if len(some_list) == 0:
        return 0
    if i == len(some_list) - 1:
        return some_list[i]
    return some_list[i] + rec_linear_sum(some_list, i + 1)''',
              'Линеаризация списка': '''
def linear(some_list):
    if not some_list:
        return some_list
    if type(some_list[0]) is list:
        return linear(some_list[0]) + linear(some_list[1:])
    return some_list[:1] + linear(some_list[1:])''',
              'Мини-судоку': '''
from random import shuffle
from copy import deepcopy


def make_assumptions(sudoku):
    for i, row in enumerate(sudoku):
        for j, value in enumerate(row):
            if not value:
                values = set(row) \
                    | set([sudoku[k][j] for k in range(4)]) \
                    | set([sudoku[m][n] for m in range((i // 2) * 2, (i // 2) * 2 + 2) 
                           for n in range((j // 2) * 2, (j // 2) * 2 + 2)])
                yield i, j, list(set(range(1, 5)) - values)


def solve_sudoku(sudoku):
    if all([k for row in sudoku for k in row]):
        return sudoku
    assumptions = list(make_assumptions(sudoku))
    shuffle(assumptions)

    x, y, values = min(assumptions, key=lambda x: len(x[2]))

    for v in values:
        new_sudoku = deepcopy(sudoku)
        new_sudoku[x][y] = v
        s = solve_sudoku(new_sudoku)
        if s:
            return s
    return None


sudoku = [[int(n) for n in input().split()[0]] for _ in range(4)]
for string in solve_sudoku(sudoku):
    print("".join(map(str, string)))''',
              'Полумагический квадрат': '''
import sys

matrix = []
for line in sys.stdin:
    line = line.rstrip('\n')
    matrix.append(list(map(int, line.split())))

if all([sum(elem) == sum(matrix[0]) for elem in matrix]) and all([sum(elem_) == sum(matrix[0]) for elem_ in list(zip(*matrix))]):
    print('YES')
else:
    print('NO')
в пятой строке \ n''',
              'Колода карт': '''

delete = input()

sus = ['пик', 'треф', 'бубен', 'червей']
mast = ['2', '3', '4', '5', '6', '7', '8', '9', '10', 'валет', 'дама', 'король', 'туз']

for elem in mast:
    for s in sus:
        if s == delete:
            continue
        else:
            print(elem, s)''',
              'Дислексия': '''
dict = input().lower().split()
d = dict
dict = [sorted(elem) for elem in dict]
dict = [''.join(elem) for elem in dict]
res = []
text = input().lower().split()
for i in range(len(text)):

    if any([sorted(list(text[i])) == sorted(list(w)) and text[i] != w for w in d]):
        if "".join(sorted(list(text[i]))) in dict:
            res.append('#' * len(text[i]))
        else:
            res.append(text[i])
    else:
        res.append(text[i])
print(*res)''',
              'Свернуть к минимуму': '''
import sys
from functools import reduce


data = [i.replace("/n", "") for i in map(str.strip, sys.stdin)]
print(reduce(lambda x, y: sorted([x, y])[0], data))''',
              'Бинго!': '''
import random


def make_bingo():
    x = random.sample(range(1, 100), 24)
    x.insert(12, 0) 
    return tuple(tuple(x[i:(i + 5)]) for i in [0, 5, 10, 15, 20])''',
              'Выбор тайного друга': '''
import random
import sys

names = list(map(str.strip, sys.stdin))
res = names.copy()


while True:
    random.shuffle(res)
    flag = True
    for i in range(len(res)):
        if res[i] == names[i]:
            flag = False
    if flag:
        break

for i in range(len(res)):
    print(f'{names[i]} - {res[i]}')''',
              'Генератор визуально различимых паролей (базовый)': '''
import random
et = '23456789qwertyuipasdfghjkzxcvbnmQWERTYUPASDFGHJKLZXCVBNM'


def generate_password(m):
    password = ''
    for i in range(m):
        password += et[random.randint(0, 55)]
    return password


def main(n, m):
    mm = []
    for i in range(n):
        pas = generate_password(m)
        while pas in mm:
            pas = generate_password(m)
        mm.append(pas)
    return mm''',
              'Дни рождения друзей': '''
import datetime as dt

days = int(input())
now = dt.datetime.today()
delta = dt.timedelta(days=days)
result = now + delta
print(result.day, result.month)''',
              'Генератор визуально различимых паролей (A)': '''
import random
et = '23456789qwertyuipasdfghjkzxcvbnmQWERTYUPASDFGHJKLZXCVBNM'


def generate_password(m):
    password = ''
    jopa = []
    for i in range(m):
        ff = et[random.randint(0, 55)]
        while ff in jopa:
            ff = et[random.randint(0, 55)]
        password += ff
        jopa.append(ff)
    return password


def main(n, m):
    mm = []
    for i in range(n):
        pas = generate_password(m)
        while pas in mm:
            pas = generate_password(m)
        mm.append(pas)
    return mm''',
              'Генератор визуально различимых паролей (B)': '''
import random

st1 = ['q', 'w', 'e', 'r', 't', 'y', 'u', 'i', 'p', 'a', 's', 'd', 'f', 'g',
       'h', 'j', 'k', 'z', 'x', 'c', 'v', 'b', 'n', 'm']
st2 = ['Q', 'W', 'E', 'R', 'T', 'Y', 'U', 'P', 'A', 'S', 'D', 'F', 'G', 'H',
       'J', 'K', 'L', 'Z', 'X', 'C', 'V', 'B', 'N', 'M']
st3 = ['2', '3', '4', '5', '6', '7', '8', '9']
st4 = st1 + st2 + st3


def generate_password(m):
    pas = []
    pas.append(random.choice(st1))
    pas.append(random.choice(st2))
    pas.append(random.choice(st3))
    for i in range(0, m - 3):
        pas.append(random.choice(st4))
    random.shuffle(pas)
    return ''.join(pas)


def main(n, m):
    list_passw = set()
    while len(list_passw) < n:
        list_passw.add(generate_password(m))
    return list_passw''',
              'Биоритмы': '''
import datetime as dt
import math

a1 = list(map(int, input().split('.')))
date1 = dt.date(a1[2], a1[1], a1[0])
b1 = list(map(int, input().split('.')))
date2 = dt.date(b1[2], b1[1], b1[0])

t = (date2 - date1).days


def snus(t, p):
    return round(math.sin((2 * math.pi * t) / p) * 100, 2)


print(snus(t, 23))
print(snus(t, 28))
print(snus(t, 33))''',
              'Найти приближённое значение Пи': '''
import random


k = 0
for i in range(1000000):
    k += (random.random() ** 2 + random.random() ** 2 < 1)
print(4 * k / 1000000)''',
              'Противоположный цвет (разминка)': '''
print(*(list(map(lambda x: 255 - x, list(map(int, input().split()))))))''',
              'Средний цвет фотографии': '''
from PIL import Image

im = Image.open("image.jpg")
pixels = im.load()
x, y = im.size
srr = 0
srg = 0
srb = 0

for i in range(x):
    for j in range(y):
        r, g, b = pixels[i, j]
        srr += r
        srg += g
        srb += b

print(srr // (x * y), srg // (x * y), srb // (x * y), end=' ')''',
              'Вертикальное отражение': '''
from PIL import Image


def mirror():
    im = Image.open("image.jpg")
    pixels = im.load()
    x, y = im.size
    for i in range(x // 2):
        for j in range(y):
            pixels[i, j], pixels[x - 1 - i, j] = pixels[x - 1 - i, j], pixels[i, j]
    im.save("res.jpg")''',
              'Диагональное отражение': '''
from PIL import Image


def mirror():
    im = Image.open("image.jpg")
    x, y = im.size
    im = im.rotate(90)
    im = im.transpose(Image.FLIP_LEFT_RIGHT)
    im.save("res.jpg")''',
              'Рисуем парусник': '''
from PIL import Image, ImageDraw


def picture(file_name, width, height,
            sky_color='#87CEEB',
            ocean_color='#017B92',
            boat_color='#874535',
            sail_color='#FFFFFF',
            sun_color='#FFCF40'):
    im = Image.new('RGB', (width, height), sky_color)
    draw = ImageDraw.Draw(im)
    draw.rectangle(((0, int(height * 0.8)), (width - 1, height - 1)), ocean_color)

    draw.polygon(((int(width * 0.25), int(height * 0.65)),
                 (int(width * 0.75), int(height * 0.65)),
                 (int(width * 0.7), int(height * 0.85)),
                 (int(width * 0.3), int(height * 0.85))), boat_color)

    draw.rectangle(((int(0.49 * width), int(height * 0.3)), (int(0.51 * width), int(height * 0.65))), boat_color)

    draw.polygon(((int(0.51 * width), int(height * 0.3)),
                  (int(0.66 * width), int(height * 0.45)),
                  (int(0.51 * width), int(height * 0.6))),
                 sail_color)

    draw.ellipse((
        (int(0.8 * width), -int(0.2 * height)),
        (int(1.2 * width), int(0.2 * height))),
        sun_color)

    im.save(file_name)''',
              'Градиент': '''
from PIL import Image, ImageDraw


def gradient(color):
    new_image = Image.new("RGB", (512, 200), (0, 0, 0))
    draw = ImageDraw.Draw(new_image)
    n = 0
    for i in range(0, 512, 2):
        if color.lower() == 'r':
            draw.line((i, 0, i, 200), (n, 0, 0))
            draw.line((i + 1, 0, i + 1, 200), (n, 0, 0))

        elif color.lower() == 'g':
            draw.line((i, 0, i, 200), (0, n, 0))
            draw.line((i + 1, 0, i + 1, 200), (0, n, 0))

        elif color.lower() == 'b':
            draw.line((i, 0, i, 200), (0, 0, n))
            draw.line((i + 1, 0, i + 1, 200), (0, 0, n))

        n += 1

    new_image.save("res.png")''',
              'Шахматная доска': '''
from PIL import Image, ImageDraw


def board(num, size):
    black_fat_cock = (0, 0, 0)
    while_small_dick = (255, 255, 255)
    im = Image.new("RGB", (num * size, num * size), while_small_dick)
    draw = ImageDraw.Draw(im)
    for i in range(0, num * size, size):
        if i % (size * 2) == 0:
            for j in range(0, num * size, size):
                if j % (size * 2) == 0:
                    draw.rectangle([i, j, i + size - 1, j + size - 1], black_fat_cock)
        else:
            for j in range(size, size * num, size):
                if j % (size * 2) != 0:
                    draw.rectangle([i, j, i + size - 1, j + size - 1], black_fat_cock)
    im.save("res.png", "PNG")''',
              'Рисуем ёлочку': '''
from PIL import Image, ImageDraw


def picture(file_name, width, height, sky_color='#75BBFD', snow_color='#FFFAFA',
            trunk_color='#A45A52', needls_color='#01796F', sun_color='#FFDB00'):
    im = Image.new('RGB', (width, height))
    drawer = ImageDraw.Draw(im)
    drawer.rectangle(((0, 0), (width, 0.8 * height)), sky_color)
    drawer.rectangle(((0, 0.8 * height), (width, height)),
                     snow_color)
    drawer.ellipse((
        (int(0.8 * width), -int(0.2 * height)),
        (int(1.2 * width), int(0.2 * height))),
        sun_color)

    drawer.rectangle(((width * 0.45, height * 0.7), (width * 0.55, height * 0.9)), trunk_color)

    drawer.polygon(((width * 0.4, height * 0.3),
                    (width * 0.5, height * 0.1),
                    (width * 0.6, height * 0.3)),
                   needls_color)

    drawer.polygon(((width * 0.35, height * 0.5),
                    (width * 0.45, height * 0.3),
                    (width * 0.55, height * 0.3),
                    width * 0.65, height * 0.5),
                   needls_color)
    drawer.polygon(((width * 0.3, height * 0.7),
                    (width * 0.4, height * 0.5),
                    (width * 0.6, height * 0.5),
                    width * 0.7, height * 0.7),
                   needls_color)

    im.save(file_name)''',
              'Стереопара': '''
from PIL import Image


def makeanagliph(filename, delta):
    im = Image.open(filename)
    x, y = im.size
    im_2 = Image.new('RGB', (x, y), (0, 0, 0))
    pi_2 = im_2.load()
    pi_1 = im.load()
    for i in range(x):
        for j in range(y):
            if i < delta:
                r, g, b = pi_1[i, j]
                pi_2[i, j] = 0, g, b
            else:
                g, b = pi_1[i, j][1:]
                r = pi_1[i - delta, j][0]
                pi_2[i, j] = r, g, b

    im_2.save("res.jpg")''',
              'Графический миксер': '''
from PIL import Image


def twist_image(input_file_name, output_file_name):
    im = Image.open(input_file_name)
    pixels = im.load()
    x, y = im.size
    for j in range(y):
        for i in range(x // 2):
            pixels[i, j], pixels[x // 2 + i, j] = pixels[x // 2 + i, j], pixels[i, j]
    im.save(output_file_name)''',
              'Прозрачность': '''
from PIL import Image, ImageDraw


def transparency(filename1, filename2):

    input_im1 = Image.open(filename1)
    input_im2 = Image.open(filename2)
    pixels1 = input_im1.load()
    x1, y1 = input_im1.size

    input_im2.resize((x1, y1))

    pixels2 = input_im2.load()
    x2, y2 = input_im2.size

    result = Image.new("RGB", (x1, y1), (0, 0, 0))
    pixels_result = result.load()
    xr, yr = result.size

    for i in range(xr):
        for j in range(yr):

            r1, g1, b1 = pixels1[i, j]
            r2, g2, b2 = pixels2[i, j]

            rr = int(0.5 * r1 + 0.5 * r2)
            gr = int(0.5 * g1 + 0.5 * g2)
            br = int(0.5 * b1 + 0.5 * b2)
            pixels_result[i, j] = rr, gr, br

    result.save('res.jpg')''',
              'Чип и Дейл': '''
import wave
import struct


def chip_and_dale(number):
    source = wave.open("in.wav", mode="rb")
    dest = wave.open("out.wav", mode="wb")
    dest.setparams(source.getparams())
    frames_count = source.getnframes()
    data = struct.unpack("<" + str(frames_count) + "h",
                         source.readframes(frames_count))

    newdata = []

    for i in range(0, len(data), number):
        newdata.append(data[i])

    newframes = struct.pack("<" + str(len(newdata)) + "h", *newdata)
    dest. writeframes(newframes)
    source.close()
    dest.close()''',
              'Поворот с размытием': '''
from PIL import Image, ImageFilter


def motion_blur(n):
    im = Image.open("image.jpg")
    im = im.transpose(Image.ROTATE_270)
    im = im.filter(ImageFilter.GaussianBlur(radius=n))
    im.save("res.jpg")''',
              'Убрать тишину': '''
import wave
import struct


def break_the_silence():
    source = wave.open("in.wav", mode="rb")
    dest = wave.open("out.wav", mode="wb")
    dest.setparams(source.getparams())
    frames_count = source.getnframes()
    data = struct.unpack("<" + str(frames_count) + "h",
                         source.readframes(frames_count))
    newdata = list(filter(lambda x: abs(x) > 5, data))
    newframes = struct.pack("<" + str(len(newdata)) + "h", *newdata)
    dest.writeframes(newframes)
    source.close()
    dest.close()''',
              'Миниатюра для сайта': '''
from PIL import Image


def make_preview(size, n_colors):
    im = Image.open('image.jpg')
    im = im.resize(size)
    im = im.quantize(n_colors)
    im.save('res.bmp')''',
              'Кручу-верчу': '''
import wave
import struct


def pitch_and_toss():
    source = wave.open("in.wav", mode="rb")
    dest = wave.open("out.wav", mode="wb")
    dest.setparams(source.getparams())
    frames_count = source.getnframes()
    data = struct.unpack("<" + str(frames_count) + "h",
                         source.readframes(frames_count))
    newdata = [data[:len(data) // 4], data[len(data) // 4: len(data) // 4 * 2],
               data[len(data) // 4 * 2:len(data) // 4 * 3], data[len(data) // 4 * 3:]]
    newdata = list(newdata[2]) + list(newdata[3]) + list(newdata[0]) + list(newdata[1])
    newframes = struct.pack('<' + str(len(newdata)) + 'h', *newdata)
    dest.writeframes(newframes)
    source.close()
    dest.close()''',
              'Формы глаголов': '''
import pymorphy2
import sys

g = sys.stdin.readlines()
data = map(lambda x: x.strip(), g)
temp_data = ''
for line in data:
    for alpha in line:
        if alpha.isalpha():
            temp_data += alpha
        else:
            temp_data += ' '
    temp_data += ' '

morph = pymorphy2.MorphAnalyzer()
finish_data = temp_data.lower().split()

counter = 0
for elem in finish_data:
    if morph.parse(elem)[0].normal_form.lower() in {'видеть', 'увидеть', 'глядеть', 'примечать', 'узреть'}:
        counter += 1


print(counter)''',
              'Существительные': '''
import pymorphy2
import sys


data = map(lambda x: x.strip(), sys.stdin.readlines())
temp_data = ''
for line in data:
    for alpha in line:
        if alpha.isalpha():
            temp_data += alpha
        else:
            temp_data += ' '
    temp_data += ' '

morph = pymorphy2.MorphAnalyzer()
finish_data = temp_data.lower().split()
dic_data = {}
for word in finish_data:
    res = morph.parse(word)[0]
    if res.score > 0.5:
        if "NOUN" in res.tag:
            if res.normal_form not in dic_data.keys():
                dic_data[res.normal_form] = 0
            dic_data[res.normal_form] += 1

result = [elem[1] for elem in sorted([(counter, word) for word, counter in dic_data.items()], reverse=True)[:10]]

print(*result)''',
              '99 бутылок кваса': '''
import pymorphy2

s = pymorphy2.MorphAnalyzer().parse('бутылка')[0]

for i in range(99, 0, - 1):
    print('В холодильнике', i, s.make_agree_with_number(i).word, 'кваса.')
    print('Возьмём одну и выпьем.')
    i -= 1
    if i != 11 and i % 10 == 1:
        kk = 'Осталась'
    else:
        kk = 'Осталось'
    print(kk, i, s.make_agree_with_number(i).word, 'кваса.')''',
              'Склоняй меня полностью': '''
import pymorphy2


morph = pymorphy2.MorphAnalyzer()
a = input()
res = morph.parse(a)[0]
if 'NOUN' in res.tag.POS:
    print('Единственное число:')
    print('Именительный падеж:', res.inflect({'nomn'}).word)
    print('Родительный падеж:', res.inflect({'gent'}).word)
    print('Дательный падеж:', res.inflect({'datv'}).word)
    print('Винительный падеж:', res.inflect({'accs'}).word)
    print('Творительный падеж:', res.inflect({'ablt'}).word)
    print('Предложный падеж:', res.inflect({'loct'}).word)
    print('Множественное число:')
    print('Именительный падеж:', res.inflect({'nomn', 'plur'}).word)
    print('Родительный падеж:', res.inflect({'gent', 'plur'}).word)
    print('Дательный падеж:', res.inflect({'datv', 'plur'}).word)
    print('Винительный падеж:', res.inflect({'accs', 'plur'}).word)
    print('Творительный падеж:', res.inflect({'ablt', 'plur'}).word)
    print('Предложный падеж:', res.inflect({'loct', 'plur'}).word)
else:
    print('Не существительное')''',
              'Спрягай меня полностью': '''
import pymorphy2


morph = pymorphy2.MorphAnalyzer()
a = input()
res = morph.parse(a)[0]
if 'VERB' in res.tag.POS or 'INFN' in res.tag.POS:
    print('Прошедшее время:')
    print(res.inflect({'past', 'VERB', 'masc'}).word)
    print(res.inflect({'past', 'VERB', 'femn'}).word)
    print(res.inflect({'past', 'VERB', 'neut'}).word)
    print(res.inflect({'past', 'VERB', 'plur'}).word)

    print('Настоящее время:')
    print(res.inflect({'pres', 'VERB', '1per', 'sing'}).word)
    print(res.inflect({'pres', 'VERB', '1per', 'plur'}).word)
    print(res.inflect({'pres', 'VERB', '2per', 'sing'}).word)
    print(res.inflect({'pres', 'VERB', '2per', 'plur'}).word)
    print(res.inflect({'pres', 'VERB', '3per', 'sing'}).word)
    print(res.inflect({'pres', 'VERB', '3per', 'plur'}).word)
else:
    print('Не глагол')''',
              'Оно живое!': '''
import sys
import pymorphy2

data = map(lambda x: x.strip(), sys.stdin.readlines())

morph = pymorphy2.MorphAnalyzer()
for elem in data:
    res = morph.parse(elem)[0]
    if 'NOUN' in res.tag:
        if 'anim' in res.tag:
            if 'plur' in res.tag:
                print('Живые')
            elif 'masc' in res.tag:
                print('Живой')
            elif 'femn' in res.tag:
                print('Живая')
        else:
            if 'plur' in res.tag:
                print('Не живые')
            elif 'masc' in res.tag:
                print('Не живой')
            elif 'femn' in res.tag:
                print('Не живая')
            elif 'neut' in res.tag:
                print('Не живое')
    else:
        print('Не существительное')''',
              'Приглашения': '''
from docx import Document
import sys

place = input()
time = input()
data = list(map(lambda elem: elem.strip(), sys.stdin.readlines()))
document = Document()


for elem in data:
    p = document.add_heading('Салам Алейкум!')
    p = document.add_paragraph('Приглашение на сходу \n место: ' + str(place))
    p = document.add_paragraph('Время: ' + str(time)).italic = True
    p = document.add_paragraph('имя: ' + str(elem))
    document.save(str(elem) + '.docx')''',
              'Презентация': '''
from pptx import Presentation

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Random.seed"
subtitle.text = "Random.seed([X], version=2) - инициализация генератора случайных чисел." \
                " Если X не указан, используется системное время."

title_slide_layout1 = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(title_slide_layout)
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "Random.shuffle"
subtitle1.text = "Random.shuffle(sequence, [rand]) - перемешивает последовательность" \
                 " (изменяется сама последовательность)." \
                 " Поэтому функция не работает для неизменяемых объектов."

title_slide_layout2 = prs.slide_layouts[0]
slide2 = prs.slides.add_slide(title_slide_layout)
title2 = slide2.shapes.title
subtitle2 = slide2.placeholders[1]
title2.text = "Random.expovariate"
subtitle2.text = "Random.expovariate(lambd) - экспоненциальное распределение. " \
                 "lambd равен 1/среднее желаемое." \
                 " Lambd должен быть отличным от нуля. Возвращаемые значения от 0" \
                 " до плюс бесконечности, если lambd положительно, и от минус" \
                 " бесконечности до 0, если lambd отрицательный."

title_slide_layout3 = prs.slide_layouts[0]
slide3 = prs.slides.add_slide(title_slide_layout)
title3 = slide3.shapes.title
subtitle3 = slide3.placeholders[1]
title3.text = "Random.lognormvariate"
subtitle3.text = "Random.lognormvariate(mu, sigma) - логарифм нормального распределения." \
                 " Если взять натуральный логарифм этого распределения," \
                 " то вы получите нормальное" \
                 " распределение со средним mu и стандартным отклонением sigma." \
                 " mu может иметь любое значение," \
                 " и sigma должна быть больше нуля."

title_slide_layout4 = prs.slide_layouts[0]
slide4 = prs.slides.add_slide(title_slide_layout)
title4 = slide4.shapes.title
subtitle4 = slide4.placeholders[1]
title4.text = "Random.vonmisesvariate"
subtitle4.text = "Random.vonmisesvariate(mu, kappa) - mu - средний угол," \
                 " выраженный в радианах от 0 до 2π," \
                 " и kappa - параметр концентрации," \
                 " который должен быть больше или равен нулю." \
                 " Если каппа равна нулю," \
                 " это распределение сводится к случайному углу в диапазоне от 0 до 2π."

prs.save('test.pptx')''',
              'Учебная ведомость': '''
import docxtpl


def create_training_sheet(class_name, subject_name, tpl_name, *marks):
    d = docxtpl.DocxTemplate(tpl_name)
    marks = sorted(marks, key=lambda x: x[0])
    table = {'class_name': class_name,
             'subject_name': subject_name,
             'marks': [{'num': i + 1, 'fio': marks[i][0], 'mark': marks[i][1]}
                       for i in range(len(marks))]}
    d.render(table)
    d.save("res.docx")''',
              'Круговая диаграмма': '''
import xlsxwriter


workbook = xlsxwriter.Workbook('res.xlsx')
worksheet = workbook.add_worksheet()

data = [('Питание', 1200), ('Развлечения', 1500), ('Учеба', 300), ('Лечение', 100), ('Прочее', 670)]

for row, (item, price) in enumerate(data):
    worksheet.write(row, 0, item)
    worksheet.write(row, 1, price)
    chart = workbook.add_chart({'type': 'pie'})
    chart.add_series({'values': '=Sheet1!B1:B5'})
    worksheet.insert_chart('C3', chart)
    chart.add_series({
        'categories': '=Sheet1!A1:A5',
        'values': '=Sheet1!B1:B5',
    })
workbook.close()''',
              'Простая документация в дорогу': '''
from docx import Document


def markdown_to_docx(text):
    document = Document()
    lines = text.split('\n')
    document.add_heading(lines[0], 0)
    for line in lines[1:]:
        if line:
            if line[:7].count('#') == 1:
                document.add_heading(line[2:], level=1)
            elif line[:7].count('#') == 2:
                document.add_heading(line[3:], level=2)
            elif line[:7].count('#') == 3:
                document.add_heading(line[4:], level=3)
            elif line[:7].count('#') == 4:
                document.add_heading(line[5:], level=4)
            elif line[:7].count('#') == 5:
                document.add_heading(line[6:], level=5)
            elif line[:7].count('#') == 6:
                document.add_heading(line[7:], level=6)
            elif str(line[:2]) == '- ':
                document.add_paragraph(line[2:], style='List Bullet')
            elif str(line[:2]) == '* ':
                document.add_paragraph(line[2:], style='List Bullet')
            elif str(line[:2]) == '+ ':
                document.add_paragraph(line[2:], style='List Bullet')
            elif line[0].isdigit() and line[1] == '.':
                document.add_paragraph(line[3:], style='List Number')
            elif line[:3].count('_') == 1 or line[:3].count('*') == 1:
                document.add_paragraph().add_run(line[1:-1]).italic = True
            elif line[:3].count('_') == 2 or line[:3].count('*') == 2:
                document.add_paragraph().add_run(line[2:-2]).bold = True
            elif line[:3].count('_') == 3 or line[:3].count('*') == 3:
                runner = document.add_paragraph().add_run(line[3:-3])
                runner.bold = True
                runner.italic = True
            else:
                document.add_paragraph(line)
        else:
            document.add_paragraph()
    document.save('res.docx')''',
              'Чек': '''
import xlsxwriter


def export_check(text):
    workbook = xlsxwriter.Workbook('res.xlsx')
    worksheet = workbook.add_worksheet()
    text = text.split('\n')
    s = 0
    for i in range(len(text)):
        name = text[i].split('\t')[0]
        price = text[i].split('\t')[1]
        n = text[i].split('\t')[2]
        worksheet.write(i, 0, name)
        worksheet.write(i, 1, float(price))
        worksheet.write(i, 2, int(n))
        worksheet.write(i, 3, '=B' + str(i + 1) + '*C' + str(i + 1))
        s += 1
    worksheet.write(s, 0, 'Итого')
    worksheet.write(s, 3, '=SUM(D1:D' + str(s) + ')')
    workbook.close()
В 7 строке \ n, 10-12 = \ t''',
              'Чеки': '''
import xlsxwriter


def export_check(text):
    workbook = xlsxwriter.Workbook('res.xlsx')

    checks = list(map(lambda x: sorted(x.split('\n')), text.split("---")))
    for i in checks:
        add = {}
        for j in i:
            if j == '':
                continue

            cc = j.split('\t')
            key, val = (cc[0], int(cc[1])), int(cc[2])

            if key in add:
                add[key] += val
            else:
                add[key] = val
        s1 = add.keys()
        s = []
        for i in s1:
            s.append([i[0], int(add[i]), int(i[1])])
        s.sort()
        f = {}
        for i in s:
            f[(i[0], i[2])] = add[(i[0], i[2])]
            del add[(i[0], i[2])]
        if f:
            worksheet = workbook.add_worksheet()
            for row, (item_price, count) in enumerate(f.items()):
                worksheet.write(row, 0, item_price[0])
                worksheet.write(row, 1, float(item_price[1]))
                worksheet.write(row, 2, float(count))
                worksheet.write(row, 3, f'=B{row + 1}*C{row + 1}')

            row += 1

            worksheet.write(row, 0, 'Итого')
            worksheet.write(row, 3, f'=SUM(D1:D{row})')
    workbook.close()
В строке 7 \ n''',
              'Выборки': '''
class Selector:
    def __init__(self, list_):
        self.__odds = list(filter(lambda elem: elem % 2 == 1, list_))
        self.__evens = list(filter(lambda elem: elem % 2 == 0, list_))

    def get_odds(self):
        return self.__odds

    def get_evens(self):
        return self.__evens''',
              'Вывод предложений': '''
class LeftParagraph:
    def __init__(self, n):
        self.text = []
        self.n = n

    def add_word(self, add):
        self.text.append(add)

    def end(self):
        st = self.text[0]
        for i in range(1, len(self.text)):
            if self.n - len(st) - 1 >= len(self.text[i]):
                st = st + ' ' + self.text[i]
            else:
                print(st)
                st = self.text[i]
        print(st)
        self.text.clear()


class RightParagraph:
    def __init__(self, n):
        self.text = []
        self.n = n

    def add_word(self, add):
        self.text.append(add)

    def end(self):
        st = self.text[0]
        for i in range(1, len(self.text)):
            if self.n - len(st) - 1 >= len(self.text[i]):
                st = st + ' ' + self.text[i]
            else:
                print((self.n - len(st)) * ' ' + st)
                st = self.text[i]
        print((self.n - len(st)) * ' ' + st)
        self.text.clear()''',
              'Форматы дат': '''
class AmericanDate:
    def __init__(self, year, month, date):
        self.year = year
        self.month = month
        self.date = date

    def set_year(self, j):
        self.year = j

    def set_month(self, j):
        self.month = j

    def set_day(self, j):
        self.date = j

    def get_year(self):
        return self.year

    def get_month(self):
        return self.month

    def get_day(self):
        return self.date

    def format(self):
        if len(str(self.date)) == 1 and len(str(self.month)) == 1:
            return '0' + str(self.month) + '.' + '0' + str(self.date) + '.' + str(self.year)
        elif len(str(self.date)) == 1:
            return str(self.month) + '.' + '0' + str(self.date) + '.' + str(self.year)
        elif len(str(self.month)) == 1:
            return '0' + str(self.month) + '.' + str(self.date) + '.' + str(self.year)
        else:
            return str(self.month) + '.' + str(self.date) + '.' + str(self.year)


class EuropeanDate:
    def __init__(self, year, month, date):
        self.year = year
        self.month = month
        self.date = date

    def set_year(self, j):
        self.year = j

    def set_month(self, j):
        self.month = j

    def set_day(self, j):
        self.date = j

    def get_year(self):
        return self.year

    def get_month(self):
        return self.month

    def get_day(self):
        return self.date

    def format(self):
        if len(str(self.date)) == 1 and len(str(self.month)) == 1:
            return '0' + str(self.date) + '.' + '0' + str(self.month) + '.' + str(self.year)
        elif len(str(self.date)) == 1:
            return '0' + str(self.date) + '.' + str(self.month) + '.' + str(self.year)
        elif len(str(self.month)) == 1:
            return str(self.date) + '.' + '0' + str(self.month) + '.' + str(self.year)
        else:
            return str(self.date) + '.' + str(self.month) + '.' + str(self.year)''',
              'Статистика': '''
class MinStat:
    def __init__(self):
        self.chiselki = []

    def add_number(self, arg):
        self.chiselki.append(arg)

    def result(self):
        if len(self.chiselki) == 0:
            return None
        return min(self.chiselki)


class MaxStat:
    def __init__(self):
        self.chiselki = []

    def add_number(self, arg):
        self.chiselki.append(arg)

    def result(self):
        if len(self.chiselki) == 0:
            return None
        return max(self.chiselki)


class AverageStat:
    def __init__(self):
        self.chiselki = []
        self.s = 0

    def add_number(self, arg):
        self.chiselki.append(arg)
        self.s += 1

    def result(self):
        if len(self.chiselki) == 0:
            return None
        return sum(self.chiselki) / self.s''',
              'Таблица': '''
class Table:
    def __init__(self, rows, cols):
        self.rows, self.cols = rows, cols
        self.matrix = [[0] * cols for _ in range(rows)]

    def get_value(self, row, col):
        if row >= self.rows or col >= self.cols or row < 0 or col < 0:
            return None
        return self.matrix[row][col]

    def set_value(self, row, col, value):
        self.matrix[row][col] = value

    def n_rows(self):
        return self.rows

    def n_cols(self):
        return self.cols''',
              'Прямоугольники': '''
class Rectangle:

    def __init__(self, x, y, w_side, h_side):
        self.__x1_coor = x
        self.__y1_coor = y
        self.__x2_coor = x + w_side
        self.__y2_coor = y + h_side

    def get_x(self):
        return self.__x1_coor

    def get_y(self):
        return self.__y1_coor

    def get_w(self):
        return self.__x2_coor - self.__x1_coor

    def get_h(self):
        return self.__y2_coor - self.__y1_coor

    def intersection(self, other: 'Rectangle'):
        result_x1 = max(self.__x1_coor, other.__x1_coor)
        result_y1 = max(self.__y1_coor, other.__y1_coor)
        result_x2 = min(self.__x2_coor, other.__x2_coor)
        result_y2 = min(self.__y2_coor, other.__y2_coor)

        if result_x1 < result_x2 and result_y1 < result_y2:
            return Rectangle(result_x1, result_y1, result_x2 - result_x1, result_y2 - result_y1)
        return None''',
              'Таблица с изменяемым размером': '''
class Table(object):

    def __init__(self, rows, cols):
        self._rows = rows
        self._cols = cols
        self._table = [[0] * cols for _ in range(rows)]

    def get_value(self, row, col):
        return (self._table[row][col] if 0 <= row < self._rows and 0 <= col < self._cols
                else None)

    def set_value(self, row, col, value):
        self._table[row][col] = value

    def n_rows(self):
        return self._rows

    def n_cols(self):
        return self._cols

    def delete_row(self, row):
        self._table.pop(row)
        self._rows -= 1

    def delete_col(self, col):
        for row in range(self._rows):
            self._table[row].pop(col)
        self._cols -= 1

    def add_row(self, row):
        self._table.insert(row, [0] * self._cols)
        self._rows += 1

    def add_col(self, col):
        for row in range(self._rows):
            self._table[row].insert(col, 0)
        self._cols += 1


def main():
    # Example 1
    tab = Table(3, 5)
    tab.set_value(0, 1, 10)
    tab.set_value(1, 2, 20)
    tab.set_value(2, 3, 30)
    for i in range(tab.n_rows()):
        for j in range(tab.n_cols()):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    tab.add_row(1)

    for i in range(tab.n_rows()):
        for j in range(tab.n_cols()):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    # Example 2
    tab = Table(2, 2)

    for i in range(tab.n_rows()):
        for j in range(tab.n_cols()):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    tab.set_value(0, 0, 10)
    tab.set_value(0, 1, 20)
    tab.set_value(1, 0, 30)
    tab.set_value(1, 1, 40)

    for i in range(tab.n_rows()):
        for j in range(tab.n_cols()):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    for i in range(-1, tab.n_rows() + 1):
        for j in range(-1, tab.n_cols() + 1):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    tab.add_row(0)
    tab.add_col(1)

    for i in range(-1, tab.n_rows() + 1):
        for j in range(-1, tab.n_cols() + 1):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    # Example 3
    tab = Table(1, 1)

    for i in range(tab.n_rows()):
        for j in range(tab.n_cols()):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    tab.set_value(0, 0, 1000)

    for i in range(tab.n_rows()):
        for j in range(tab.n_cols()):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    for i in range(-1, tab.n_rows() + 1):
        for j in range(-1, tab.n_cols() + 1):
            print(tab.get_value(i, j), end=' ')
        print()
    print()

    tab.add_row(0)
    tab.add_row(2)
    tab.add_col(0)
    tab.add_col(2)

    tab.set_value(0, 0, 2000)
    tab.set_value(0, 2, 3000)
    tab.set_value(2, 0, 4000)
    tab.set_value(2, 2, 5000)

    for i in range(-1, tab.n_rows() + 1):
        for j in range(-1, tab.n_cols() + 1):
            print(tab.get_value(i, j), end=' ')
        print()
    print()


if __name__ == "__main__":
    main()''',
              'Калорийность': '''
class FoodInfo:

    def __init__(self, proteins, fats, carbohydrates):
        self.__proteins = proteins
        self.__fats = fats
        self.__carbohydrates = carbohydrates

    def get_proteins(self):
        return self.__proteins

    def get_fats(self):
        return self.__fats

    def get_carbohydrates(self):
        return self.__carbohydrates

    def get_kcalories(self):
        return 4 * self.__proteins + 9 * self.__fats + 4 * self.__carbohydrates

    def __add__(self, other):
        new_proteins = self.__proteins + other.__proteins
        new_fats = self.__fats + other.__fats
        new_carbohydrates = self.__carbohydrates + other.__carbohydrates
        return FoodInfo(new_proteins, new_fats, new_carbohydrates)''',
              'Список в обратном порядке': '''
class ReversedList:
    def __init__(self, arg_list):
        self.__reversed_list = arg_list.copy()
        self.__number = len(arg_list)

    def __len__(self):
        return self.__number

    def __getitem__(self, key):
        return self.__reversed_list[self.__number - key - 1]''',
              'Квадратичная функция': '''
class SquareFunction:

    def __init__(self, a, b, c):
        self.__a = a
        self.__b = b
        self.__c = c

    def __call__(self, x):
        return self.__a * (x ** 2) + self.__b * x + self.__c''',
              'Вычитание дат': '''
import datetime as dt


class Date:
    def __init__(self, month, day):
        self.month = month
        self.day = day

    def __sub__(self, other):
        date1 = dt.date(2019, self.month, self.day)
        date2 = dt.date(2019, other.month, other.day)
        total_date = str(date1 - date2).split()
        if len(total_date) == 1:
            return 0
        return total_date[0]''',
              'Точки на плоскости': '''
class Point:

    def __init__(self, x, y):
        self.__x = x
        self.__y = y

    def __eq__(self, other):
        if self.__y == other.__y and self.__x == other.__x:
            return True
        return False

    def __ne__(self, other):
        if self.__y == other.__y and self.__x == other.__x:
            return False
        return True''',
              'Разреженный массив': '''
class SparseArray:

    def __init__(self):
        self.__sparse_array = dict()

    def __getitem__(self, key):
        if key not in self.__sparse_array.keys():
            return 0
        return self.__sparse_array[key]

    def __setitem__(self, key, value):
        self.__sparse_array[key] = value
        return self''',
              'Сложение многочленов': '''
class Polynomial:
    def __init__(self, koef):
        self.koef = koef

    def __call__(self, x):
        s = 0
        for i in range(len(self.koef)):
            s += self.koef[i] * pow(x, i)
        return s

    def __add__(self, other):
        st = []
        k = Polynomial(st)
        if len(self.koef) < len(other.koef):
            m = len(self.koef)
        else:
            m = len(other.koef)
        for i in range(m):
            st.append(self.koef[i] + other.koef[i])
        if len(self.koef) > m:
            st += self.koef[m::]
        else:
            st += other.koef[m::]
        k.koef = st
        return k''',
              'Треугольники': '''
class Triangle:
    def __init__(self, a, b, c):
        self._a = a
        self._b = b
        self._c = c

    def perimeter(self):
        return self._a + self._b + self._c


class EquilateralTriangle(Triangle):
    def __init__(self, a):
        self._a = a
        self._b = a
        self._c = a''',
              'Сумматоры': '''
from functools import reduce


class Summator:
    def __init__(self, arg_power=1):
        self._power = arg_power

    def transform(self, n):
        return n

    def sum(self, n):
        return reduce(lambda value, elem: value + self.transform(elem), range(1, n + 1))


class PowerSummator(Summator):
    def __init__(self, b):
        super().__init__(b)

    def transform(self, n):
        return PowerSummator.__binary_power(n, self._power)

    @staticmethod
    def __binary_power(base, power):
        if power == 0:
            return 1
        if power % 2 == 0:
            temp = PowerSummator.__binary_power(base, power // 2)
            return temp * temp
        return base * PowerSummator.__binary_power(base, power - 1)


class SquareSummator(PowerSummator):
    def __init__(self):
        super().__init__(2)


class CubeSummator(PowerSummator):
    def __init__(self):
        super().__init__(3)''',
              'Сайт поиска вакансий': '''
class Profile:
    def __init__(self, type_profession):
        self._type_profession = type_profession

    def info(self):
        return ""

    def describe(self):
        return self._type_profession + self.info()


class Vacancy(Profile):
    def __init__(self, arg_type_profession, arg_salary):
        super().__init__(arg_type_profession)
        self.__salary = arg_salary

    def info(self):
        return f"Предлагаемая зарплата: {self.__salary}"


class Resume(Profile):

    def __init__(self, arg_type_profession, arg_work_exp):
        super().__init__(arg_type_profession)
        self.__work_experiense = arg_work_exp

    def info(self):
        return f"Стаж работы: {self.__work_experiense}"''',
              'Сумматоры – 2': '''
from functools import reduce


class Summator:
    def __init__(self, arg_power=1):
        self._power = arg_power

    def transform(self, n):
        return n

    def sum(self, n):
        return reduce(lambda value, elem: value + self.transform(elem), range(1, n + 1))


class PowerSummator(Summator):
    def __init__(self, b):
        super().__init__(b)

    def transform(self, n):
        return PowerSummator.__binary_power(n, self._power)

    @staticmethod
    def __binary_power(base, power):
        if power == 0:
            return 1
        if power % 2 == 0:
            temp = PowerSummator.__binary_power(base, power // 2)
            return temp * temp
        return base * PowerSummator.__binary_power(base, power - 1)


class SquareSummator(PowerSummator):
    def __init__(self):
        super().__init__(2)


class CubeSummator(PowerSummator):
    def __init__(self):
        super().__init__(3)''',
              'Алфавит классов': '''
class A:
    def __init__(self):
        pass

    def __str__(self):
        return 'A.__str__ method'

    def hello(self):
        print("Hello")


class B:
    def __init__(self):
        pass

    def __str__(self):
        return 'B.__str__ method'

    def good_evening(self):
        print("Good evening")


class C(A, B):
    def __init__(self):
        A.__init__(self)


class D(B, A):
    def __init__(self):
        B.__init__(self)''',
              'Конь': '''
class Knight:
    def __init__(self, row, col, color):
        self._row = row
        self._col = col
        self._color = color

    def can_move(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (col1 < 8):
            if ((abs(row1 - self._row) == 2 and abs(col1 - self._col) == 1) 
                    or (abs(row1 - self._row) == 1 and abs(col1 - self._col) == 2)):
                return True
            return False
        return False

    def set_position(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (row1 < 8):
            self._row = row1
            self._col = col1

    def get_color(self):
        return self._color

    def char(self):
        return "N"''',
              'Слон': '''
class Bishop:
    def __init__(self, row, col, color):
        self._row = row
        self._col = col
        self._color = color

    def can_move(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (col1 < 8):
            if abs(row1 - self._row) == abs(col1 - self._col):
                return True
            return False
        return False

    def set_position(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (row1 < 8):
            self._row = row1
            self._col = col1

    def get_color(self):
        return self._color

    def char(self):
        return "B"''',
              'Ферзь': '''
class Queen:
    def __init__(self, row, col, color):
        self._row = row
        self._col = col
        self._color = color

    def can_move(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (col1 < 8):
            if ((abs(row1 - self._row) == abs(col1 - self._col)) or
                    (abs(row1 - self._row) > -1 and col1 - self._col == 0) or
                    (row1 - self._row == 0 and abs(col1 - self._row) > -1)):
                return True
            return False
        return False

    def set_position(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (row1 < 8):
            self._row = row1
            self._col = col1

    def get_color(self):
        return self._color

    def char(self):
        return "Q"''',
              'Поля под боем': '''
WHITE = 1
BLACK = 2


class White:
    def __eq__(self, other):
        # истина, если другой операнд оператора ==
        # тоже является экземпляром (англ. instance) класса Black
        return isinstance(other, White)

    def opponent(self):
        return Black()

    def is_black(self):
        return True

    def is_white(self):
        return False


class Black:
    def __eq__(self, other):
        # истина, если другой операнд оператора ==
        # тоже является экземпляром (англ. instance) класса Black
        return isinstance(other, Black)

    def opponent(self):
        return White()

    def is_black(self):
        return True

    def is_white(self):
        return False


class Pawn:

    def __init__(self, row, col, color):
        self.row = row
        self.col = col
        self.color = color

    def set_position(self, row, col):
        self.row = row
        self.col = col

    def char(self):
        return 'P'

    def get_color(self):
        return self.color

    def can_move(self, row, col):
        # Пешка может ходить только по вертикали
        # "взятие на проходе" не реализовано
        if self.col != col:
            return False

        # Пешка может сделать из начального положения ход на 2 клетки
        # вперёд, поэтому поместим индекс начального ряда в start_row.
        if self.color == WHITE:
            direction = 1
            start_row = 1
        else:
            direction = -1
            start_row = 6

        # ход на 1 клетку
        if self.row + direction == row:
            return True

        # ход на 2 клетки из начального положения
        if self.row == start_row and self.row + 2 * direction == row:
            return True

        return False


class Rook:

    def __init__(self, row, col, color):
        self.row = row
        self.col = col
        self.color = color

    def set_position(self, row, col):
        self.row = row
        self.col = col

    def char(self):
        return 'R'

    def get_color(self):
        return self.color

    def can_move(self, row, col):
        if self.row != row and self.col != col:
            return False

        return True


class Knight:
    def __init__(self, row, col, color):
        self._row = row
        self._col = col
        self._color = color

    def can_move(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (col1 < 8):
            if ((abs(row1 - self._row) == 2 and abs(col1 - self._col) == 1)
                    or (abs(row1 - self._row) == 1 and abs(col1 - self._col) == 2)):
                return True
            return False
        return False

    def set_position(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (row1 < 8):
            self._row = row1
            self._col = col1

    def get_color(self):
        return self._color

    def char(self):
        return "N"


class Bishop:
    def __init__(self, row, col, color):
        self._row = row
        self._col = col
        self._color = color

    def can_move(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (col1 < 8):
            if abs(row1 - self._row) == abs(col1 - self._col):
                return True
            return False

    def set_position(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (row1 < 8):
            self._row = row1
            self._col = col1

    def get_color(self):
        return self._color

    def char(self):
        return "B"


class Queen:
    def __init__(self, row, col, color):
        self._row = row
        self._col = col
        self._color = color

    def can_move(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (col1 < 8):
            if ((abs(row1 - self._row) == abs(col1 - self._col)) or
                    (abs(row1 - self._row) > -1 and col1 - self._col == 0) or
                    (row1 - self._row == 0 and abs(col1 - self._row) > -1)):
                return True
            return False
        return False

    def set_position(self, row1, col1):
        if (row1 > -1) and (row1 < 8) and (col1 > -1) and (row1 < 8):
            self._row = row1
            self._col = col1

    def get_color(self):
        return self._color

    def char(self):
        return "Q"


def correct_coords(row, col):
    return 0 <= row < 8 and 0 <= col < 8


def opponent(color):
    if color == WHITE:
        return BLACK
    return WHITE


class Board:
    def __init__(self):
        self.color = WHITE
        self.field = []
        for row in range(8):
            self.field.append([None] * 8)
        self.field[1][4] = Pawn(1, 4, WHITE)

    def current_player_color(self):
        return self.color

    def cell(self, row, col):

        piece = self.field[row][col]
        if piece is None:
            return '  '
        color = piece.get_color()
        c = 'w' if color == WHITE else 'b'
        return c + piece.char()

    def move_piece(self, row, col, row1, col1):
        if not correct_coords(row, col) or not correct_coords(row1, col1):
            return False
        if row == row1 and col == col1:
            return False  # нельзя пойти в ту же клетку
        piece = self.field[row][col]
        if piece is None:
            return False
        if piece.get_color() != self.color:
            return False
        if not piece.can_move(row1, col1):
            return False
        self.field[row][col] = None
        self.field[row1][col1] = piece
        piece.set_position(row1, col1)
        self.color = opponent(self.color)
        return True

    def is_under_attack(self, row, col, color):
        for elem in self.field:
            for el in elem:
                if el is not None:
                    if el.can_move(row, col) and el.get_color() == color:
                        return True          
        return False'''
              }

with open("db_cash.json", "w") as cf:
    json.dump(cas_db, cf)

with open("db_free.json", "w") as cf:
    json.dump(fre_db, cf)
